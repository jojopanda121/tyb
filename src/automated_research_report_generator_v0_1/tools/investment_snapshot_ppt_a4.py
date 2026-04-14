from __future__ import annotations

import calendar
import re
import math
from dataclasses import dataclass
from pathlib import Path
from typing import Type

from crewai.tools import BaseTool
from pydantic import BaseModel, Field
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


"""横向 A4 单页投资快照 PPT 导出工具"""


FALLBACK_TEXT = "缺乏信息"
PRIMARY_FONT = "Microsoft YaHei"
TITLE_FONT = "SimHei"
SLIDE_WIDTH = Inches(11.69)
SLIDE_HEIGHT = Inches(8.27)
SAFE_LEFT = Inches(0.82)
SAFE_TOP = Inches(0.42)
SAFE_BOTTOM = Inches(0.44)
SAFE_WIDTH = int(SLIDE_WIDTH - (SAFE_LEFT * 2))
TITLE_BLOCK_WIDTH = Inches(3.05)
TITLE_RULE_TOP = Inches(0.86)
BANNER_TOP = Inches(1.10)
TOP_PANEL_TOP = Inches(1.74)
TOP_PANEL_GUTTER = Inches(0.22)
TOP_PANEL_WIDTH = int((SAFE_WIDTH - TOP_PANEL_GUTTER) / 2)
TOP_PANEL_HEIGHT = Inches(2.56)
MID_PANEL_GAP = Inches(0.30)
RISK_PANEL_GAP = Inches(0.22)
BACKGROUND_COLOR = RGBColor(255, 255, 255)
BRICK_RED = RGBColor(154, 52, 28)
GOLD = RGBColor(198, 149, 82)
LIGHT_BEIGE = RGBColor(239, 232, 223)
SOFT_BEIGE = RGBColor(243, 238, 232)
TABLE_ROW_FILL_PRIMARY = RGBColor(248, 244, 238)
TABLE_ROW_FILL_SECONDARY = RGBColor(241, 235, 228)
TEXT_DARK = RGBColor(33, 28, 24)
DASH_GOLD = RGBColor(208, 177, 130)
REQUIRED_FINANCIAL_LABELS = (
    "营业收入",
    "营业收入增长率",
    "毛利率",
    "净利率",
    "净利润增长率",
    "总资产",
    "资产负债率",
    "ROE",
)
DIRECT_FINANCIAL_LABEL_ALIASES = {
    "营业收入": ("营业收入",),
    "营业收入增长率": ("营业收入增长率", "收入增速(%)", "营业收入增速", "收入增长率"),
    "毛利率": ("毛利率", "销售毛利率(%)", "销售毛利率"),
    "净利率": ("净利率", "销售净利率(%)", "销售净利率"),
    "净利润增长率": ("净利润增长率", "净利润增速(%)", "净利润增长", "净利增长率"),
    "总资产": ("总资产", "资产总计"),
    "资产负债率": ("资产负债率", "负债率"),
    "ROE": ("ROE",),
}
SUPPORTING_FINANCIAL_LABEL_ALIASES = {
    "净利润": ("净利润",),
    "股东权益": ("股东权益", "所有者权益", "权益总额"),
}


class InvestmentSnapshotFinancialRow(BaseModel):  # 目的：定义渲染用财务行结构；功能：承接解析或推导后的单个指标行；实现逻辑：使用标签加期间值列表统一传给表格绘制层；可调参数：label、values；默认参数及原因：字段均必填，避免表格生成时再猜测结构。
    label: str = Field(..., description="财务指标名称，例如营业收入")
    values: list[str] = Field(..., description="与展示期间对齐的期间值列表")


class InvestmentSnapshotOverviewProductItem(BaseModel):  # 目的：定义公司概况中的产品项结构；功能：约束产品名称与一句话说明；实现逻辑：将产品名和描述拆开，便于 PPT 中分别应用粗体和常规样式；可调参数：name、description；默认参数及原因：字段均必填，避免正文拼接时丢失标题层次。
    name: str = Field(..., description="产品名称，会在 PPT 中加粗显示")
    description: str = Field(..., description="约 20-30 字的产品简介与用途")


class InvestmentSnapshotTitledItem(BaseModel):  # 目的：定义亮点或风险项结构；功能：约束短标题和展开说明；实现逻辑：以标题和正文分栏建模，便于富文本段落渲染；可调参数：title、detail；默认参数及原因：字段均必填，防止 agent 退化成无结构长句。
    title: str = Field(..., description="短标题，会在 PPT 中加粗显示")
    detail: str = Field(..., description="对标题的进一步阐述")


class InvestmentSnapshotPptInput(BaseModel):  # 目的：定义单页 PPT 工具入参；功能：把叙事内容与财务原文固定为结构化字段；实现逻辑：通过明确 schema 限制 agent 输出形态，减少版式和解析歧义；可调参数：路径、标题、概况、财务原文、亮点和风险；默认参数及原因：关键字段全部必填，确保工具可以独立完成导出。
    pptx_path: str = Field(..., description="输出 PPTX 文件路径")
    slide_title: str = Field(..., description="页面主标题，固定建议为 投资要点速览")
    positioning_line: str = Field(..., description="公司定位短句，会与公司名拼接为红色 banner")
    overview_summary: str = Field(..., description="约 90-120 字的公司业务与吸引力摘要")
    overview_product_items: list[InvestmentSnapshotOverviewProductItem] = Field(..., description="公司概况中的产品项，固定建议 3 条")
    financial_source_markdown: str = Field(..., description="conduct_financial_analysis 的完整 Markdown 原文")
    highlight_items: list[InvestmentSnapshotTitledItem] = Field(..., description="投资亮点条目，固定建议 3-5 条")
    risk_items: list[InvestmentSnapshotTitledItem] = Field(..., description="投资风险条目，固定建议 1-2 条")


@dataclass
class ParsedFinancialTable:  # 目的：承接 Markdown 财务表解析结果；功能：统一保存原始期间、展示期间、选中下标、单位和指标映射；实现逻辑：解析层一次提取，后续映射和绘图直接复用；可调参数：periods、display_periods、selected_indexes、unit_note、row_lookup；默认参数及原因：不设置运行时默认值，避免 silently 产生错误口径。
    periods: list[str]
    display_periods: list[str]
    selected_indexes: list[int]
    unit_note: str
    row_lookup: dict[str, list[str]]


class InvestmentSnapshotPptTool(BaseTool):  # 目的：定义横向 A4 单页投资快照 PPT 导出工具；功能：根据结构化叙事和财务 Markdown 原文直接绘制真实 PPTX；实现逻辑：agent 只压缩叙事，工具负责财务解析、映射推导、版式绘制和文件导出；可调参数：全部 schema 字段；默认参数及原因：工具名固定且返回结果作为答案，便于任务直接产出文件消息。
    name: str = "investment_snapshot_ppt_tool"
    description: str = (
        "Create a single-slide investment snapshot PowerPoint (.pptx) with a fixed "
        "A4 investment-committee layout using structured narrative content and the "
        "full financial analysis markdown."
    )
    args_schema: Type[BaseModel] = InvestmentSnapshotPptInput

    def _run(  # 目的：执行 PPT 导出主流程；功能：规范化叙事内容、解析财务原文并生成单页 PPTX；实现逻辑：先整理输入，再由工具内部确定性构建财务表，最后统一绘制整页；可调参数：全部工具入参；默认参数及原因：不增加隐式默认字段，避免 agent 输出和成品之间出现不可见偏差。
        self,
        pptx_path: str,
        slide_title: str,
        positioning_line: str,
        overview_summary: str,
        overview_product_items: list[InvestmentSnapshotOverviewProductItem],
        financial_source_markdown: str,
        highlight_items: list[InvestmentSnapshotTitledItem],
        risk_items: list[InvestmentSnapshotTitledItem],
    ) -> str:
        out_file = Path(pptx_path).expanduser().resolve()
        out_file.parent.mkdir(parents=True, exist_ok=True)

        parsed_financial_table = self._parse_financial_markdown(financial_source_markdown)
        presentation = self._build_presentation(
            company_name=self._derive_company_name_from_output_path(out_file),
            slide_title=self._normalize_llm_text(slide_title, fallback_text="投资要点速览"),
            positioning_line=self._normalize_llm_text(positioning_line, fallback_text=FALLBACK_TEXT),
            overview_summary=self._normalize_llm_text(overview_summary, fallback_text=FALLBACK_TEXT),
            overview_product_items=self._normalize_overview_product_items(overview_product_items, 3, 3),
            financial_table=parsed_financial_table,
            financial_rows=self._build_snapshot_financial_rows(parsed_financial_table),
            highlight_items=self._normalize_titled_items(highlight_items, 3, 5, "投资亮点"),
            risk_items=self._normalize_titled_items(risk_items, 1, 2, "投资风险"),
        )
        presentation.save(str(out_file))
        return f"PPT created successfully at: {out_file}"

    def _build_presentation(  # 目的：创建演示文稿对象；功能：设置横向 A4 画布并绘制整页内容；实现逻辑：统一创建 blank slide 后依次绘制标题、概况、财务、亮点和风险区域；可调参数：标准化后的内容字段；默认参数及原因：单页 blank layout 能最大限度避免模板主题干扰。
        self,
        company_name: str,
        slide_title: str,
        positioning_line: str,
        overview_summary: str,
        overview_product_items: list[InvestmentSnapshotOverviewProductItem],
        financial_table: ParsedFinancialTable,
        financial_rows: list[InvestmentSnapshotFinancialRow],
        highlight_items: list[InvestmentSnapshotTitledItem],
        risk_items: list[InvestmentSnapshotTitledItem],
    ) -> Presentation:
        presentation = Presentation()
        presentation.slide_width = SLIDE_WIDTH
        presentation.slide_height = SLIDE_HEIGHT

        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = BACKGROUND_COLOR

        self._add_title_area(slide, company_name, slide_title, positioning_line)
        self._add_overview_panel(slide, overview_summary, overview_product_items)
        self._add_financial_panel(slide, financial_table, financial_rows)
        risk_top = self._add_highlights_panel(slide, highlight_items) + RISK_PANEL_GAP
        self._add_risk_panel(slide, risk_items, risk_top)
        return presentation

    def _derive_company_name_from_output_path(self, pptx_path: Path) -> str:  # 目的：从固定输出文件名推导公司名；功能：避免额外增加 company_name 入参；实现逻辑：优先剥离 _investment_snapshot 后缀，失败时退回 stem；可调参数：输出路径；默认参数及原因：当前任务链已经固定输出命名，沿用约定最简单稳妥。
        stem = pptx_path.stem.strip()
        suffix = "_investment_snapshot"
        if stem.endswith(suffix):
            candidate = stem[: -len(suffix)].strip()
            if candidate:
                return candidate
        return stem or "未知公司"

    def _normalize_overview_product_items(self, items: list[InvestmentSnapshotOverviewProductItem] | list[dict[str, object]], min_items: int, max_items: int) -> list[InvestmentSnapshotOverviewProductItem]:  # 目的：规范化产品项列表；功能：保证公司概况区固定为 3 条“产品名 + 描述”；实现逻辑：清洗空项并在缺失时补充占位项，但不对 LLM 正文做字符截断；可调参数：items、min_items、max_items；默认参数及原因：固定 3 条能稳定 A4 版式且符合用户要求。
        normalized_items: list[InvestmentSnapshotOverviewProductItem] = []
        for raw_item in items[:max_items]:
            name, description = self._extract_named_item_payload(raw_item)
            if not name and not description:
                continue
            normalized_items.append(
                InvestmentSnapshotOverviewProductItem(
                    name=self._normalize_llm_text(name, fallback_text="核心产品"),
                    description=self._normalize_llm_text(description, fallback_text=FALLBACK_TEXT),
                )
            )
        while len(normalized_items) < min_items:
            normalized_items.append(InvestmentSnapshotOverviewProductItem(name="核心产品", description=FALLBACK_TEXT))
        return normalized_items

    def _normalize_titled_items(self, items: list[InvestmentSnapshotTitledItem] | list[dict[str, object]], min_items: int, max_items: int, fallback_title: str) -> list[InvestmentSnapshotTitledItem]:  # 目的：规范化亮点或风险条目；功能：限制条数并补足缺失项；实现逻辑：先抽取结构化字段，再做空值兜底和占位补全，但不对 LLM 正文做字符截断；可调参数：items、min_items、max_items、fallback_title；默认参数及原因：最少条数用于避免版块空白，最大条数用于保证版式稳定。
        normalized_items: list[InvestmentSnapshotTitledItem] = []
        for raw_item in items[:max_items]:
            title, detail = self._extract_titled_item_payload(raw_item)
            if not title and not detail:
                continue
            normalized_items.append(
                InvestmentSnapshotTitledItem(
                    title=self._normalize_llm_text(title, fallback_text=fallback_title),
                    detail=self._normalize_llm_text(detail, fallback_text=FALLBACK_TEXT),
                )
            )
        while len(normalized_items) < min_items:
            normalized_items.append(InvestmentSnapshotTitledItem(title=fallback_title, detail=FALLBACK_TEXT))
        return normalized_items

    def _normalize_llm_text(self, text: str | None, fallback_text: str) -> str:  # 目的：规范化 LLM 叙事文本；功能：只做空白清洗和缺失兜底，不对内容长度做字符截断；实现逻辑：统一压缩多余空白并在空值时回填指定默认值；可调参数：text、fallback_text；默认参数及原因：坚持由提示词约束内容长度，工具不再二次裁剪。
        normalized_text = " ".join(str(text or "").split())
        return normalized_text or fallback_text

    def _normalize_financial_text(self, text: str | None, fallback_text: str) -> str:  # 目的：规范化财务表文本；功能：只做空白清洗和缺失兜底，不截断任何财务数值或期间值；实现逻辑：统一压缩多余空白并在空值时回填指定默认值；可调参数：text、fallback_text；默认参数及原因：财务表数据精确性优先，绝不因版式原因裁掉字符。
        normalized_text = " ".join(str(text or "").split())
        return normalized_text or fallback_text

    def _extract_named_item_payload(self, item: InvestmentSnapshotOverviewProductItem | dict[str, object]) -> tuple[str, str]:  # 目的：兼容不同来源的产品项对象；功能：同时支持 Pydantic 模型和 CrewAI 传入的 dict；实现逻辑：统一抽取 name 与 description 两个字段；可调参数：item；默认参数及原因：无法识别时返回空字符串，便于上层统一补位而不是直接报错。
        if isinstance(item, InvestmentSnapshotOverviewProductItem):
            return str(item.name).strip(), str(item.description).strip()
        if isinstance(item, dict):
            return str(item.get("name", "")).strip(), str(item.get("description", "")).strip()
        return "", ""

    def _extract_titled_item_payload(self, item: InvestmentSnapshotTitledItem | dict[str, object]) -> tuple[str, str]:  # 目的：兼容不同来源的亮点或风险对象；功能：同时支持 Pydantic 模型和 CrewAI 传入的 dict；实现逻辑：统一抽取 title 与 detail 两个字段；可调参数：item；默认参数及原因：无法识别时返回空字符串，交给上层占位逻辑处理更稳妥。
        if isinstance(item, InvestmentSnapshotTitledItem):
            return str(item.title).strip(), str(item.detail).strip()
        if isinstance(item, dict):
            return str(item.get("title", "")).strip(), str(item.get("detail", "")).strip()
        return "", ""

    def _parse_financial_markdown(self, financial_source_markdown: str) -> ParsedFinancialTable:  # 目的：解析财务分析 Markdown 原文；功能：从核心财务表中提取期间、单位和指标行，并只保留最新 3 个展示期间；实现逻辑：定位标题区块后抓取 Markdown 表格，再按列名拆解数据；可调参数：financial_source_markdown；默认参数及原因：展示固定 3 期是横向 A4 下最稳妥的密度选择。
        markdown = financial_source_markdown.strip()
        if not markdown:
            raise ValueError("financial_source_markdown must not be empty.")
        lines = markdown.splitlines()
        heading_index = next((index for index, line in enumerate(lines) if line.strip().startswith("## 1. 核心财务数据总表")), -1)
        if heading_index < 0:
            raise ValueError("Failed to find '## 1. 核心财务数据总表' in financial_source_markdown.")

        table_lines: list[str] = []
        found_table = False
        for line in lines[heading_index + 1 :]:
            stripped = line.strip()
            if stripped.startswith("## ") and found_table:
                break
            if stripped.startswith("|"):
                table_lines.append(stripped)
                found_table = True
                continue
            if found_table:
                break
        if len(table_lines) < 3:
            raise ValueError("Failed to parse the markdown financial table from financial_source_markdown.")

        header_cells = self._parse_markdown_table_row(table_lines[0])
        period_column_end = self._find_period_column_end(header_cells)
        raw_periods = [cell for cell in header_cells[1:period_column_end] if cell]
        if len(raw_periods) < 3:
            raise ValueError("The financial markdown table must contain at least 3 periods.")

        selected_indexes = list(range(max(0, len(raw_periods) - 3), len(raw_periods)))
        row_lookup: dict[str, list[str]] = {}
        note_candidates: list[str] = []
        for line in table_lines[2:]:
            if self._is_markdown_separator_row(line):
                continue
            cells = self._parse_markdown_table_row(line)
            if len(cells) < 1 + len(raw_periods):
                continue
            label = cells[0].strip()
            if not label:
                continue
            row_lookup[label] = [
                self._normalize_financial_text(value, fallback_text=FALLBACK_TEXT)
                for value in cells[1 : 1 + len(raw_periods)]
            ]
            if len(cells) > 1 + len(raw_periods):
                note_candidates.extend(cells[1 + len(raw_periods) :])

        return ParsedFinancialTable(
            periods=raw_periods,
            display_periods=[self._format_period_label(raw_periods[index]) for index in selected_indexes],
            selected_indexes=selected_indexes,
            unit_note=self._extract_financial_meta_note(markdown, note_candidates),
            row_lookup=row_lookup,
        )

    def _parse_markdown_table_row(self, line: str) -> list[str]:  # 目的：解析单行 Markdown 表格；功能：把 `|` 分隔的文本行拆成单元格列表；实现逻辑：去掉首尾分隔符后按竖线切分并逐格去空白；可调参数：line；默认参数及原因：只处理标准 Markdown 表格，保持实现简单可控。
        return [cell.strip() for cell in line.strip().strip("|").split("|")]

    def _is_markdown_separator_row(self, line: str) -> bool:  # 目的：识别 Markdown 表头下的分隔行；功能：避免把 `| :--- |` 这一行误当成数据；实现逻辑：检查整行是否仅由竖线、冒号、横线和空格组成；可调参数：line；默认参数及原因：只需支持仓库现有的标准 Markdown 表格写法。
        stripped = line.strip()
        return bool(stripped) and all(char in "|:- " for char in stripped)

    def _find_period_column_end(self, header_cells: list[str]) -> int:  # 目的：确定财务期间列在表头中的结束位置；功能：排除“数据来源”“备注”等说明列；实现逻辑：从第二列开始扫描，一旦遇到说明列标题即停止；可调参数：header_cells；默认参数及原因：兼容 5 列以上的研究表格结构而不硬编码列数。
        meta_patterns = ("数据来源", "备注", "定义", "口径", "页码")
        for index, header in enumerate(header_cells[1:], start=1):
            if any(pattern in header for pattern in meta_patterns):
                return index
        return len(header_cells)

    def _format_period_label(self, period: str) -> str:  # 目的：整理期间展示文本；功能：把常见日期口径统一转成 YYYY/MM/DD，减少表头视觉噪声；实现逻辑：优先识别完整日期、年末日期、括号日期和常见区间月份，无法确定具体日时才回退原文；可调参数：period；默认参数及原因：财务表头以最短可读日期为主，不再使用“2024年（12/31）”这类冗长混排。
        normalized_period = " ".join(period.split())
        year_match = re.fullmatch(r"(\d{4})[-/年](?:12[-/月]?31日?)", normalized_period)
        if year_match:
            return f"{year_match.group(1)}/12/31"
        parenthesized_month_day_match = re.fullmatch(r"(\d{4})年[（(](\d{1,2})/(\d{1,2})[)）]", normalized_period)
        if parenthesized_month_day_match:
            return f"{parenthesized_month_day_match.group(1)}/{int(parenthesized_month_day_match.group(2)):02d}/{int(parenthesized_month_day_match.group(3)):02d}"
        month_span_match = re.fullmatch(r"(\d{4})年(?:\d{1,2})[-~至](\d{1,2})月", normalized_period)
        if month_span_match:
            year = int(month_span_match.group(1))
            month = int(month_span_match.group(2))
            day = calendar.monthrange(year, month)[1]
            return f"{year}/{month:02d}/{day:02d}"
        explicit_date_patterns = (
            r"(\d{4})年(\d{1,2})月(\d{1,2})日?",
            r"(\d{4})/(\d{1,2})/(\d{1,2})",
            r"(\d{4})-(\d{1,2})-(\d{1,2})",
        )
        for explicit_date_pattern in explicit_date_patterns:
            full_date_match = re.search(explicit_date_pattern, normalized_period)
            if full_date_match:
                return f"{full_date_match.group(1)}/{int(full_date_match.group(2)):02d}/{int(full_date_match.group(3)):02d}"
        quarter_match = re.search(r"(\d{4}).*Q([1-4])", normalized_period, flags=re.IGNORECASE)
        if quarter_match:
            year = int(quarter_match.group(1))
            month = int(quarter_match.group(2)) * 3
            day = calendar.monthrange(year, month)[1]
            return f"{year}/{month:02d}/{day:02d}"
        return normalized_period

    def _extract_financial_meta_note(self, markdown: str, note_candidates: list[str]) -> str:  # 目的：提取财务表“币种 + 单位”口径；功能：兼容 LLM 可能输出的“币种为”“单位为”“金额单位”“货币单位”等多种写法；实现逻辑：先从全文和备注候选中分别提取币种与单位，再做合并、拆分和缺失兜底；可调参数：markdown、note_candidates；默认参数及原因：字段沿用 unit_note 命名以减少改动，但实际展示内容升级为币种与单位组合。
        currency_value = ""
        unit_value = ""
        for candidate in [markdown, *note_candidates]:
            for normalized_candidate_line in self._iterate_financial_meta_lines(candidate):
                if not currency_value:
                    currency_value = self._extract_currency_value(normalized_candidate_line)
                if not unit_value:
                    unit_value = self._extract_unit_value(normalized_candidate_line)
                if currency_value and unit_value:
                    break
            if currency_value and unit_value:
                break
        currency_value, unit_value = self._split_combined_currency_unit(currency_value, unit_value)
        if currency_value and unit_value:
            return f"{currency_value}，{unit_value}"
        if currency_value:
            return f"{currency_value}，单位未披露"
        if unit_value:
            return f"币种未披露，{unit_value}"
        return "币种未披露，单位未披露"

    def _extract_currency_value(self, text: str) -> str:  # 目的：提取财务表币种信息；功能：从正文或备注中识别“币种为人民币”“货币：美元”等不稳定写法；实现逻辑：依次按多组正则扫描，命中后统一清洗并返回；可调参数：text；默认参数及原因：未命中时返回空字符串，便于上层继续从其他候选补足。
        currency_patterns = (
            r"(?:币种|币别|货币)\s*(?:为|是|[:：])\s*([^\|\n。；;，,（）()]+)",
            r"(?:以|按)\s*([^\|\n。；;，,（）()]+)\s*计价",
            r"[(（]\s*(人民币|美元|港元|欧元|英镑|日元|新台币)\s*[)）]",
        )
        for pattern in currency_patterns:
            match = re.search(pattern, text, flags=re.IGNORECASE)
            if match:
                return self._clean_financial_meta_value(match.group(1))
        return ""

    def _extract_unit_value(self, text: str) -> str:  # 目的：提取财务表单位信息；功能：兼容“单位为千元”“金额单位：人民币百万元”“货币单位：亿元”等多种 LLM 表达；实现逻辑：先尝试更具体的金额单位类模式，再尝试一般单位模式，最后统一清洗结果；可调参数：text；默认参数及原因：未命中时返回空字符串，避免误把整句说明塞进表头。
        unit_patterns = (
            r"(?:金额单位|货币单位|财务单位)\s*(?:为|是|[:：])\s*([^\|\n。；;，,]+)",
            r"(?:单位|口径)\s*(?:为|是|[:：])\s*([^\|\n。；;，,]+)",
        )
        for pattern in unit_patterns:
            match = re.search(pattern, text, flags=re.IGNORECASE)
            if match:
                return self._clean_financial_meta_value(match.group(1))
        return ""

    def _split_combined_currency_unit(self, currency_value: str, unit_value: str) -> tuple[str, str]:  # 目的：拆分被写在同一字段里的币种和单位；功能：处理“人民币百万元”“亿元（人民币）”“美元千元”等组合写法；实现逻辑：优先识别单位前缀中的币种，再识别单位括号里的币种，最后返回拆分后的两部分；可调参数：currency_value、unit_value；默认参数及原因：只做有限规则拆分，不做猜测性映射，避免引入口径错误。
        known_currencies = ("人民币", "美元", "港元", "欧元", "英镑", "日元", "新台币")
        normalized_currency = self._clean_financial_meta_value(currency_value)
        normalized_unit = self._clean_financial_meta_value(unit_value)
        if not normalized_unit:
            return normalized_currency, normalized_unit
        for known_currency in known_currencies:
            if not normalized_currency and normalized_unit.startswith(known_currency):
                normalized_currency = known_currency
                normalized_unit = normalized_unit[len(known_currency) :].strip(" /、，,")
                break
        if not normalized_currency:
            bracket_match = re.search(r"[（(]\s*(人民币|美元|港元|欧元|英镑|日元|新台币)\s*[)）]", unit_value)
            if bracket_match:
                normalized_currency = bracket_match.group(1)
                normalized_unit = re.sub(r"[（(]\s*(人民币|美元|港元|欧元|英镑|日元|新台币)\s*[)）]", "", normalized_unit).strip(" /、，,")
        return normalized_currency, normalized_unit

    def _iterate_financial_meta_lines(self, candidate: str) -> list[str]:  # 目的：把财务元数据候选拆成逐行文本；功能：避免“单位”行和“会计准则”行在空白压缩后被误拼接；实现逻辑：保留原始换行逐行清洗，并过滤掉空白行；可调参数：candidate；默认参数及原因：币种、单位、会计准则通常分行出现，逐行匹配是最小且最稳的修复。
        return [
            normalized_line
            for raw_line in str(candidate).splitlines()
            if (normalized_line := " ".join(str(raw_line).split()))
        ]

    def _clean_financial_meta_value(self, value: str) -> str:  # 目的：清洗财务元数据文本；功能：去掉括号补充说明、字段前缀和多余标点，保留表头可读的核心币种或单位；实现逻辑：先压缩空白，再裁掉“除百分比外”等尾注与包裹标点，最后返回简洁结果；可调参数：value；默认参数及原因：尽量保留真实口径，同时避免把整句说明抬进表头导致视觉噪声。
        normalized_value = " ".join(str(value or "").split())
        normalized_value = re.sub(r"^(?:币种|币别|货币|金额单位|货币单位|财务单位|单位|口径)\s*(?:为|是|[:：])\s*", "", normalized_value)
        normalized_value = re.sub(r"[（(][^）)]*(?:除百分比[^）)]*|百分比外|比率外|except[^\])）]*)[）)]", "", normalized_value, flags=re.IGNORECASE)
        normalized_value = re.sub(r"[（(][A-Za-z]{2,8}[)）]", "", normalized_value)
        normalized_value = normalized_value.strip("：:；;，,。 ")
        return normalized_value

    def _build_snapshot_financial_rows(self, parsed_table: ParsedFinancialTable) -> list[InvestmentSnapshotFinancialRow]:  # 目的：构建快照页财务表 8 行数据；功能：按固定指标顺序完成别名映射、缺失补位和有限代码推导；实现逻辑：优先取显式披露值，再按规则用上游原始行计算缺失比率；可调参数：parsed_table；默认参数及原因：未命中的指标统一填“缺乏信息”，避免伪造数据。
        row_lookup = parsed_table.row_lookup
        revenue_series = self._pick_series_by_aliases(row_lookup, DIRECT_FINANCIAL_LABEL_ALIASES["营业收入"])
        net_profit_series = self._pick_series_by_aliases(row_lookup, SUPPORTING_FINANCIAL_LABEL_ALIASES["净利润"])
        total_assets_series = self._pick_series_by_aliases(row_lookup, DIRECT_FINANCIAL_LABEL_ALIASES["总资产"])
        equity_series = self._pick_series_by_aliases(row_lookup, SUPPORTING_FINANCIAL_LABEL_ALIASES["股东权益"])
        result_rows: list[InvestmentSnapshotFinancialRow] = []
        for label in REQUIRED_FINANCIAL_LABELS:
            display_values = self._select_display_values(self._pick_series_by_aliases(row_lookup, DIRECT_FINANCIAL_LABEL_ALIASES[label]), parsed_table.selected_indexes)
            if not self._has_meaningful_values(display_values):
                if label == "营业收入增长率":
                    display_values = self._select_display_values(self._compute_growth_series(revenue_series), parsed_table.selected_indexes)
                elif label == "净利率":
                    display_values = self._select_display_values(self._compute_ratio_series(net_profit_series, revenue_series), parsed_table.selected_indexes)
                elif label == "净利润增长率":
                    display_values = self._select_display_values(self._compute_growth_series(net_profit_series), parsed_table.selected_indexes)
                elif label == "资产负债率":
                    display_values = self._select_display_values(self._compute_debt_ratio_series(total_assets_series, equity_series), parsed_table.selected_indexes)
            result_rows.append(
                InvestmentSnapshotFinancialRow(
                    label=label,
                    values=[
                        self._normalize_financial_text(value, fallback_text=FALLBACK_TEXT)
                        for value in display_values
                    ],
                )
            )
        return result_rows

    def _pick_series_by_aliases(self, row_lookup: dict[str, list[str]], aliases: tuple[str, ...]) -> list[str]:  # 目的：按别名集合选取完整指标序列；功能：从解析后的行映射中找到首个命中的指标行；实现逻辑：按照别名优先级顺序逐一查找；可调参数：row_lookup、aliases；默认参数及原因：未命中时返回空列表，便于后续统一补缺和推导。
        for alias in aliases:
            if alias in row_lookup:
                return list(row_lookup[alias])
        return []

    def _select_display_values(self, series: list[str], selected_indexes: list[int]) -> list[str]:  # 目的：从完整期间序列中截取展示窗口；功能：把全量历史值压缩为最新 3 期显示值；实现逻辑：按 selected_indexes 顺序取值，缺失位补“缺乏信息”；可调参数：series、selected_indexes；默认参数及原因：保证所有行都和表头期间严格对齐。
        if not series:
            return [FALLBACK_TEXT for _ in selected_indexes]
        return [series[index] if index < len(series) else FALLBACK_TEXT for index in selected_indexes]

    def _has_meaningful_values(self, values: list[str]) -> bool:  # 目的：判断一行展示值是否已经有效命中；功能：区分“全为空占位”与“存在真实内容”；实现逻辑：只要任意单元格不是缺失占位就视为命中；可调参数：values；默认参数及原因：显式文本如“无法直接计算”也应视为有效披露。
        return any(value.strip() and value.strip() != FALLBACK_TEXT for value in values)

    def _compute_growth_series(self, series: list[str]) -> list[str]:  # 目的：计算增长率序列；功能：在缺少显式增长率行时基于上一期数值计算变化率；实现逻辑：使用 `(本期-上期)/abs(上期)`，首期或分母缺失时填缺失；可调参数：series；默认参数及原因：对亏损扩大或收窄的净利润更能直观表达改善方向。
        numeric_values = [self._parse_numeric_value(value) for value in series]
        growth_values: list[str] = []
        for index, current_value in enumerate(numeric_values):
            if index == 0:
                growth_values.append(FALLBACK_TEXT)
                continue
            previous_value = numeric_values[index - 1]
            if current_value is None or previous_value is None or previous_value == 0:
                growth_values.append(FALLBACK_TEXT)
                continue
            growth_values.append(self._format_percent_value((current_value - previous_value) / abs(previous_value) * 100))
        return growth_values

    def _compute_ratio_series(self, numerator_series: list[str], denominator_series: list[str]) -> list[str]:  # 目的：计算比率序列；功能：在缺少显式净利率行时用净利润除以营业收入计算；实现逻辑：分子分母同时可解析且分母非零时输出百分比；可调参数：numerator_series、denominator_series；默认参数及原因：仅做直接财务比率，不引入更复杂的会计推导。
        ratio_values: list[str] = []
        for index in range(max(len(numerator_series), len(denominator_series))):
            numerator = self._parse_numeric_value(numerator_series[index]) if index < len(numerator_series) else None
            denominator = self._parse_numeric_value(denominator_series[index]) if index < len(denominator_series) else None
            if numerator is None or denominator in (None, 0):
                ratio_values.append(FALLBACK_TEXT)
                continue
            ratio_values.append(self._format_percent_value(numerator / denominator * 100))
        return ratio_values

    def _compute_debt_ratio_series(self, total_assets_series: list[str], equity_series: list[str]) -> list[str]:  # 目的：计算资产负债率序列；功能：在缺少显式负债率时基于总资产和股东权益推导；实现逻辑：按 `(总资产-股东权益)/总资产` 计算，分母无效时填缺失；可调参数：total_assets_series、equity_series；默认参数及原因：只使用财务表中已存在的基本平衡关系，不引入额外假设。
        debt_ratio_values: list[str] = []
        for index in range(max(len(total_assets_series), len(equity_series))):
            total_assets = self._parse_numeric_value(total_assets_series[index]) if index < len(total_assets_series) else None
            equity = self._parse_numeric_value(equity_series[index]) if index < len(equity_series) else None
            if total_assets in (None, 0) or equity is None:
                debt_ratio_values.append(FALLBACK_TEXT)
                continue
            debt_ratio_values.append(self._format_percent_value((1 - equity / total_assets) * 100))
        return debt_ratio_values

    def _parse_numeric_value(self, value: str) -> float | None:  # 目的：把表格文本转成可计算数值；功能：兼容逗号、括号负数和百分号等常见财务格式；实现逻辑：先排除明显不可计算文本，再提取首个数字并还原负号；可调参数：value；默认参数及原因：不可计算时返回 None，避免推导逻辑误把文字当成 0。
        normalized_value = str(value).strip()
        if not normalized_value or normalized_value in {FALLBACK_TEXT, "-", "—", "N/A", "n/a"}:
            return None
        if any(flag in normalized_value for flag in ("无法", "未披露", "不可", "不适用")):
            return None
        is_negative = normalized_value.startswith("(") and normalized_value.endswith(")")
        numeric_source = normalized_value.strip("()").replace(",", "")
        match = re.search(r"[-+]?\d+(?:\.\d+)?", numeric_source)
        if match is None:
            return None
        parsed_value = float(match.group())
        if is_negative and parsed_value > 0:
            return -parsed_value
        return parsed_value

    def _format_percent_value(self, value: float) -> str:  # 目的：格式化百分比结果；功能：统一把推导出的浮点数转成一位小数百分比文本；实现逻辑：先做四舍五入再拼接 `%`；可调参数：value；默认参数及原因：一位小数兼顾精度和表格可读性。
        return f"{round(value, 1):.1f}%"

    def _truncate_text(self, text: str, max_chars: int) -> str:  # 目的：限制单段文字长度；功能：在不改变事实前提下避免文本框内容过长溢出；实现逻辑：按字符数硬截断并在尾部追加省略号；可调参数：text、max_chars；默认参数及原因：不同区域由调用方传入不同阈值，便于局部调优。
        normalized_text = " ".join(text.split())
        if len(normalized_text) <= max_chars:
            return normalized_text
        return normalized_text[: max_chars - 1].rstrip() + "…"

    def _fit_overview_font_size(self, summary: str, product_items: list[InvestmentSnapshotOverviewProductItem]) -> float:  # 目的：估算公司概况区字号；功能：根据摘要和产品项总长度决定是否缩小字体；实现逻辑：用摘要与 3 个产品项的总字符数做启发式档位缩放；可调参数：summary、product_items；默认参数及原因：优先保留 10.8pt，在内容偏多时逐档下调。
        total_chars = len(summary) + sum(len(item.name) + len(item.description) for item in product_items)
        if total_chars > 300:
            return 10.6
        if total_chars > 250:
            return 10.8
        return 11.0

    def _fit_titled_items_font_size(self, items: list[InvestmentSnapshotTitledItem], base_size: float, min_size: float, chars_per_line: float | None = None, max_content_height: int | None = None, space_after: int = 3) -> float:  # 目的：估算亮点或风险区字号；功能：根据条目数量、详情长度和可用高度粗调字体；实现逻辑：先做字符量启发式收缩，再在给定高度内逐档下调直到估算内容可放下；可调参数：items、base_size、min_size、chars_per_line、max_content_height、space_after；默认参数及原因：先保留 11pt 正文字号，仅在页面高度不足时才收缩。
        total_chars = sum(len(item.title) + len(item.detail) for item in items)
        font_size = base_size
        if total_chars > 420:
            font_size -= 0.2
        if total_chars > 520:
            font_size -= 0.3
        if total_chars > 620:
            font_size -= 0.5
        if chars_per_line is not None and max_content_height is not None:
            while font_size > min_size:
                content_height = self._estimate_titled_items_content_height(items, font_size, chars_per_line, space_after)
                if content_height <= max_content_height:
                    break
                font_size = round(font_size - 0.2, 1)
        return max(min_size, font_size)

    def _estimate_titled_items_content_height(self, items: list[InvestmentSnapshotTitledItem], font_size: float, chars_per_line: float, space_after: int) -> int:  # 目的：估算亮点或风险正文的内容高度；功能：在不依赖 PowerPoint 自动排版回读的前提下，用启发式方法估算富文本段落所需高度；实现逻辑：对每条“标题 + 正文”按加权字符数估算换行数，再累计行高和段后距；可调参数：items、font_size、chars_per_line、space_after；默认参数及原因：标题按更高权重计入，是因为粗体短标题会比正文更早触发行宽占用。
        total_height = 0
        safe_chars_per_line = max(1.0, chars_per_line)
        line_height = float(font_size) * 1.42
        for index, item in enumerate(items):
            weighted_chars = len(item.title) * 1.45 + len(item.detail)
            estimated_lines = max(1, math.ceil(weighted_chars / safe_chars_per_line))
            total_height += int(Pt(line_height * estimated_lines))
            if index < len(items) - 1:
                total_height += int(Pt(space_after))
        return total_height

    def _add_title_area(self, slide, company_name: str, slide_title: str, positioning_line: str) -> None:  # 目的：绘制页面顶部标题区；功能：输出主标题、金色分割线和红色 banner；实现逻辑：在安全边距内使用文本框与色块组合构建页首层次；可调参数：slide、company_name、slide_title、positioning_line；默认参数及原因：标题区尺寸固定，便于下方四区布局稳定对齐。
        self._add_textbox(
            slide=slide,
            left=SAFE_LEFT,
            top=SAFE_TOP,
            width=SAFE_WIDTH,
            height=Inches(0.40),
            text=f"投资要点速览——{company_name}",
            font_size=20,
            bold=True,
            font_name=TITLE_FONT,
            color=TEXT_DARK,
            align=PP_ALIGN.LEFT,
            margin_left=2,
            margin_right=2,
            margin_top=0,
            margin_bottom=0,
        )

        line_shape = slide.shapes.add_textbox(SAFE_LEFT, TITLE_RULE_TOP, SAFE_WIDTH, Inches(0.05))
        line_shape.fill.solid()
        line_shape.fill.fore_color.rgb = GOLD
        line_shape.line.fill.background()

        self._add_textbox(
            slide=slide,
            left=SAFE_LEFT,
            top=BANNER_TOP,
            width=SAFE_WIDTH,
            height=Inches(0.36),
            text=positioning_line,
            font_size=14,
            bold=True,
            font_name=TITLE_FONT,
            color=RGBColor(255, 255, 255),
            fill_color=BRICK_RED,
            line_color=BRICK_RED,
            align=PP_ALIGN.LEFT,
            margin_left=12,
            margin_right=12,
            margin_top=2,
            margin_bottom=0,
        )

    def _add_overview_panel(self, slide, overview_summary: str, overview_product_items: list[InvestmentSnapshotOverviewProductItem]) -> None:  # 目的：绘制公司概况面板；功能：呈现一段摘要和 3 条产品简介；实现逻辑：先画虚线框和标签，再用富文本段落输出摘要与“产品名加粗 + 描述”结构；可调参数：slide、overview_summary、overview_product_items；默认参数及原因：左上区固定尺寸，便于与财务区并排。
        panel_left = SAFE_LEFT
        panel_top = TOP_PANEL_TOP
        panel_width = TOP_PANEL_WIDTH
        panel_height = TOP_PANEL_HEIGHT

        self._add_panel_frame(slide, panel_left, panel_top, panel_width, panel_height)
        self._add_section_label(slide, "公司概况", panel_left + int((panel_width - Inches(1.92)) / 2), panel_top - Inches(0.22), Inches(1.92))

        textbox = self._create_textbox_container(
            slide,
            panel_left + Inches(0.14),
            panel_top + Inches(0.16),
            panel_width - Inches(0.28),
            panel_height - Inches(0.24),
        )
        text_frame = textbox.text_frame
        body_font_size = self._fit_overview_font_size(overview_summary, overview_product_items)
        self._add_plain_paragraph(text_frame, overview_summary, body_font_size, TEXT_DARK, True, 6)
        for item in overview_product_items:
            self._add_named_item_paragraph(
                text_frame=text_frame,
                item=item,
                font_size=max(10.8, body_font_size),
                title_color=BRICK_RED,
                detail_color=TEXT_DARK,
                space_after=3,
            )

    def _add_financial_panel(self, slide, financial_table: ParsedFinancialTable, financial_rows: list[InvestmentSnapshotFinancialRow]) -> None:  # 目的：绘制财务数据面板；功能：呈现最新 3 期的 8 行固定指标；实现逻辑：标题标签、表头文本框和单张 8 行表格分层组合；可调参数：slide、financial_table、financial_rows；默认参数及原因：表头独立于表格，能在 A4 中更灵活地控制列宽和字号。
        panel_left = SAFE_LEFT + TOP_PANEL_WIDTH + TOP_PANEL_GUTTER
        panel_top = TOP_PANEL_TOP
        panel_width = TOP_PANEL_WIDTH
        panel_height = TOP_PANEL_HEIGHT
        period_count = len(financial_table.display_periods)

        self._add_panel_frame(slide, panel_left, panel_top, panel_width, panel_height, show_fill=False)
        self._add_section_label(slide, "财务数据", panel_left + int((panel_width - Inches(1.92)) / 2), panel_top - Inches(0.22), Inches(1.92))

        content_left = panel_left + Inches(0.10)
        content_width = panel_width - Inches(0.20)
        first_column_width = Inches(1.66)
        other_column_width = int((content_width - first_column_width) / period_count)
        column_widths = [first_column_width] + [other_column_width] * period_count
        header_top = panel_top + Inches(0.16)
        header_height = Inches(0.36)
        header_labels = [f"项目（{financial_table.unit_note}）"] + financial_table.display_periods

        running_left = content_left
        for index, header_text in enumerate(header_labels):
            self._add_textbox(
                slide=slide,
                left=running_left,
                top=header_top,
                width=column_widths[index],
                height=header_height,
                text=header_text,
                font_size=10.2,
                bold=True,
                font_name=PRIMARY_FONT,
                color=TEXT_DARK,
                align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE,
                margin_left=2,
                margin_right=2,
                margin_top=0,
                margin_bottom=0,
            )
            running_left += column_widths[index]

        for rule_top in (header_top - Inches(0.02), header_top + header_height + Inches(0.01)):
            rule_shape = slide.shapes.add_textbox(content_left, rule_top, content_width, Inches(0.02))
            rule_shape.fill.solid()
            rule_shape.fill.fore_color.rgb = BRICK_RED
            rule_shape.line.fill.background()

        table_shape = slide.shapes.add_table(
            rows=len(financial_rows),
            cols=1 + period_count,
            left=content_left,
            top=panel_top + Inches(0.52),
            width=content_width,
            height=panel_height - Inches(0.60),
        )
        table = table_shape.table
        table.first_row = False
        table.horz_banding = False
        table.vert_banding = False
        for column_index, column_width in enumerate(column_widths):
            table.columns[column_index].width = column_width

        row_height = int((panel_height - Inches(0.60)) / max(1, len(financial_rows)))
        for row_index, row_payload in enumerate(financial_rows):
            row_fill_color = TABLE_ROW_FILL_PRIMARY if row_index % 2 == 0 else TABLE_ROW_FILL_SECONDARY
            table.rows[row_index].height = row_height
            self._set_table_cell_text(table.cell(row_index, 0), row_payload.label, 10.2, True, PP_ALIGN.LEFT, row_fill_color)
            for value_index, value in enumerate(row_payload.values, start=1):
                self._set_table_cell_text(table.cell(row_index, value_index), value, 10.2, False, PP_ALIGN.RIGHT, row_fill_color)

    def _add_highlights_panel(self, slide, highlight_items: list[InvestmentSnapshotTitledItem]) -> int:  # 目的：绘制投资亮点面板；功能：呈现 3-5 条“短标题加粗 + 进一步阐述”条目；实现逻辑：先根据内容长度估算所需高度，再绘制整宽虚线框和富文本正文；可调参数：slide、highlight_items；默认参数及原因：根据文字多少动态收口高度，避免固定高度带来的大块空白。
        panel_left = SAFE_LEFT
        panel_top = TOP_PANEL_TOP + TOP_PANEL_HEIGHT + MID_PANEL_GAP
        panel_width = SAFE_WIDTH
        max_panel_height = Inches(2.14)
        min_panel_height = Inches(1.36)
        font_size = self._fit_titled_items_font_size(
            highlight_items,
            11.0,
            10.4,
            chars_per_line=86,
            max_content_height=int(max_panel_height - Inches(0.24)),
            space_after=3,
        )
        content_height = self._estimate_titled_items_content_height(highlight_items, font_size, 86, 3)
        panel_height = max(min_panel_height, min(max_panel_height, content_height + int(Inches(0.24))))

        self._add_panel_frame(slide, panel_left, panel_top, panel_width, panel_height)
        self._add_section_label(slide, "投资亮点", panel_left + int((panel_width - Inches(1.98)) / 2), panel_top - Inches(0.22), Inches(1.98))

        textbox = self._create_textbox_container(
            slide,
            panel_left + Inches(0.14),
            panel_top + Inches(0.15),
            panel_width - Inches(0.28),
            panel_height - Inches(0.22),
        )
        text_frame = textbox.text_frame
        for index, item in enumerate(highlight_items):
            self._add_titled_item_paragraph(text_frame, item, font_size, BRICK_RED, TEXT_DARK, index == 0, 3)
        return panel_top + panel_height

    def _add_risk_panel(self, slide, risk_items: list[InvestmentSnapshotTitledItem], panel_top) -> None:  # 目的：绘制底部风险条带；功能：用左侧标签和右侧富文本正文突出 1-2 条核心风险；实现逻辑：先按文字量估算条带高度，再在统一宽度下绘制虚线外框与内部标签区、正文区；可调参数：slide、risk_items、panel_top；默认参数及原因：风险区跟随亮点区动态下移，保持整体留白和区块节奏稳定。
        panel_left = SAFE_LEFT
        panel_width = SAFE_WIDTH
        max_panel_height = Inches(1.26)
        min_panel_height = Inches(0.92)
        label_width = Inches(1.54)
        max_body_height = int(min(max_panel_height - Inches(0.16), SLIDE_HEIGHT - SAFE_BOTTOM - panel_top - Inches(0.08)))
        font_size = self._fit_titled_items_font_size(
            risk_items,
            11.0,
            10.4,
            chars_per_line=70,
            max_content_height=max_body_height,
            space_after=2,
        )
        content_height = self._estimate_titled_items_content_height(risk_items, font_size, 70, 2)
        panel_height = max(min_panel_height, min(max_panel_height, content_height + int(Inches(0.16))))

        self._add_panel_frame(slide, panel_left, panel_top, panel_width, panel_height)

        self._add_textbox(
            slide=slide,
            left=panel_left + Inches(0.02),
            top=panel_top + Inches(0.02),
            width=label_width - Inches(0.02),
            height=panel_height - Inches(0.04),
            text="投资风险",
            font_size=14,
            bold=True,
            font_name=TITLE_FONT,
            color=BRICK_RED,
            fill_color=LIGHT_BEIGE,
            line_color=LIGHT_BEIGE,
            align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
            margin_left=2,
            margin_right=2,
            margin_top=0,
            margin_bottom=0,
        )
        self._add_textbox(
            slide=slide,
            left=panel_left + label_width,
            top=panel_top + Inches(0.02),
            width=panel_width - label_width - Inches(0.02),
            height=panel_height - Inches(0.04),
            text="",
            font_size=11,
            bold=False,
            fill_color=SOFT_BEIGE,
            line_color=SOFT_BEIGE,
        )

        textbox = self._create_textbox_container(
            slide,
            panel_left + label_width + Inches(0.08),
            panel_top + Inches(0.08),
            panel_width - label_width - Inches(0.16),
            panel_height - Inches(0.16),
        )
        text_frame = textbox.text_frame
        for index, item in enumerate(risk_items):
            self._add_titled_item_paragraph(text_frame, item, font_size, BRICK_RED, TEXT_DARK, index == 0, 2)

    def _add_panel_frame(self, slide, left, top, width, height, show_fill: bool = True) -> None:  # 目的：绘制虚线内容框；功能：统一概况、财务和亮点区的边框风格；实现逻辑：使用无圆角文本框作为容器边框；可调参数：位置尺寸和是否填充；默认参数及原因：默认浅底色加金色虚线，能较接近用户给出的参考页气质。
        frame = slide.shapes.add_textbox(left, top, width, height)
        if show_fill:
            frame.fill.solid()
            frame.fill.fore_color.rgb = RGBColor(255, 255, 255)
        else:
            frame.fill.background()
        frame.line.color.rgb = DASH_GOLD
        frame.line.width = Pt(1.2)
        frame.line.dash_style = MSO_LINE_DASH_STYLE.DASH

    def _add_section_label(self, slide, text: str, left, top, width) -> None:  # 目的：绘制分区标签；功能：为每个内容面板提供居中的标题签；实现逻辑：使用浅米色文本框压在边框上方；可调参数：文字和位置；默认参数及原因：统一使用 13pt 粗体以兼顾 A4 下的层次感与空间效率。
        self._add_textbox(
            slide=slide,
            left=left,
            top=top,
            width=width,
            height=Inches(0.30),
            text=text,
            font_size=14,
            bold=True,
            font_name=TITLE_FONT,
            color=BRICK_RED,
            fill_color=LIGHT_BEIGE,
            line_color=LIGHT_BEIGE,
            align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
            margin_left=0,
            margin_right=0,
            margin_top=0,
            margin_bottom=0,
        )

    def _create_textbox_container(self, slide, left, top, width, height, anchor: MSO_ANCHOR = MSO_ANCHOR.MIDDLE):  # 目的：创建富文本正文容器；功能：统一正文文本框的换行、垂直对齐和零边框配置；实现逻辑：返回一个透明文本框，供概况、亮点和风险区继续追加段落；可调参数：位置、尺寸和垂直对齐方式；默认参数及原因：默认使用上下居中，和当前页面对齐规范保持一致。
        textbox = slide.shapes.add_textbox(left, top, width, height)
        textbox.fill.background()
        textbox.line.fill.background()
        text_frame = textbox.text_frame
        text_frame.clear()
        text_frame.word_wrap = True
        text_frame.vertical_anchor = anchor
        text_frame.margin_left = Pt(0)
        text_frame.margin_right = Pt(0)
        text_frame.margin_top = Pt(0)
        text_frame.margin_bottom = Pt(0)
        return textbox

    def _add_plain_paragraph(self, text_frame, text: str, font_size: float, color: RGBColor, is_first: bool, space_after: int) -> None:  # 目的：写入普通正文段落；功能：为概况摘要等纯文本内容统一设置字号和段后距；实现逻辑：按是否首段选择复用首个段落或新建段落，再写入单个 run；可调参数：text_frame、text、font_size、color、is_first、space_after；默认参数及原因：首段复用可避免空段落残留。
        paragraph = text_frame.paragraphs[0] if is_first else text_frame.add_paragraph()
        paragraph.alignment = PP_ALIGN.LEFT
        paragraph.space_after = Pt(space_after)
        run = paragraph.add_run()
        run.text = text
        self._apply_run_font(run, font_size, False, PRIMARY_FONT, color)

    def _add_named_item_paragraph(self, text_frame, item: InvestmentSnapshotOverviewProductItem, font_size: float, title_color: RGBColor, detail_color: RGBColor, space_after: int) -> None:  # 目的：写入“产品名加粗 + 描述”段落；功能：让公司概况中的产品层次更接近投委会阅读习惯；实现逻辑：先写粗体名称和冒号，再写常规描述 run；可调参数：text_frame、item、字号和颜色；默认参数及原因：使用独立 run 才能在测试里准确断言粗体存在。
        paragraph = text_frame.add_paragraph()
        paragraph.alignment = PP_ALIGN.LEFT
        paragraph.space_after = Pt(space_after)
        title_run = paragraph.add_run()
        title_run.text = f"{item.name}："
        self._apply_run_font(title_run, font_size, True, PRIMARY_FONT, title_color)
        detail_run = paragraph.add_run()
        detail_run.text = item.description
        self._apply_run_font(detail_run, font_size, False, PRIMARY_FONT, detail_color)

    def _add_titled_item_paragraph(self, text_frame, item: InvestmentSnapshotTitledItem, font_size: float, title_color: RGBColor, detail_color: RGBColor, is_first: bool, space_after: int) -> None:  # 目的：写入“短标题加粗 + 说明句”段落；功能：用于亮点和风险区的富文本条目渲染；实现逻辑：标题和正文分别作为不同 run 写入，同段完成强弱层次；可调参数：text_frame、item、字号、颜色、首段标记和段后距；默认参数及原因：标题单独粗体可让高密度内容仍保持扫读效率。
        paragraph = text_frame.paragraphs[0] if is_first else text_frame.add_paragraph()
        paragraph.alignment = PP_ALIGN.LEFT
        paragraph.space_after = Pt(space_after)
        title_run = paragraph.add_run()
        title_run.text = f"{item.title}："
        self._apply_run_font(title_run, font_size, True, PRIMARY_FONT, title_color)
        detail_run = paragraph.add_run()
        detail_run.text = item.detail
        self._apply_run_font(detail_run, font_size, False, PRIMARY_FONT, detail_color)

    def _add_textbox(self, slide, left, top, width, height, text: str, font_size: float, bold: bool, font_name: str = PRIMARY_FONT, color: RGBColor = TEXT_DARK, fill_color: RGBColor | None = None, line_color: RGBColor | None = None, align: PP_ALIGN = PP_ALIGN.LEFT, anchor: MSO_ANCHOR = MSO_ANCHOR.MIDDLE, margin_left: int = 6, margin_right: int = 6, margin_top: int = 4, margin_bottom: int = 4):  # 目的：封装单段文本框创建；功能：统一文本框的字体、填充、边框和内边距设置；实现逻辑：把常用文本框参数收敛到一个方法里减少重复；可调参数：位置、文本、样式与内边距；默认参数及原因：默认采用上下居中，更符合当前整页版式对齐要求。
        textbox = slide.shapes.add_textbox(left, top, width, height)
        if fill_color is None:
            textbox.fill.background()
        else:
            textbox.fill.solid()
            textbox.fill.fore_color.rgb = fill_color
        if line_color is None:
            textbox.line.fill.background()
        else:
            textbox.line.color.rgb = line_color
            textbox.line.width = Pt(0.8)

        text_frame = textbox.text_frame
        text_frame.clear()
        text_frame.word_wrap = True
        text_frame.vertical_anchor = anchor
        text_frame.margin_left = Pt(margin_left)
        text_frame.margin_right = Pt(margin_right)
        text_frame.margin_top = Pt(margin_top)
        text_frame.margin_bottom = Pt(margin_bottom)

        paragraph = text_frame.paragraphs[0]
        paragraph.alignment = align
        run = paragraph.add_run()
        run.text = text
        self._apply_run_font(run, font_size, bold, font_name, color)
        return textbox

    def _set_table_cell_text(self, cell, text: str, font_size: float, bold: bool, align: PP_ALIGN, fill_color: RGBColor = TABLE_ROW_FILL_PRIMARY) -> None:  # 目的：统一设置表格单元格文本；功能：控制字体、对齐、填充和边距，降低默认主题干扰；实现逻辑：逐格写入单个 run，并允许调用方按行传入底色以增强可读性；可调参数：cell、text、font_size、bold、align、fill_color；默认参数及原因：默认使用浅暖底色而非纯白，以提升财务表横向扫读体验。
        cell.text = ""
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill_color
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        if align == PP_ALIGN.RIGHT:
            cell.margin_left = Pt(1.2)
            cell.margin_right = Pt(3.6)
        elif align == PP_ALIGN.LEFT:
            cell.margin_left = Pt(3.6)
            cell.margin_right = Pt(1.2)
        else:
            cell.margin_left = Pt(2.0)
            cell.margin_right = Pt(2.0)
        cell.margin_top = Pt(1)
        cell.margin_bottom = Pt(1)

        text_frame = cell.text_frame
        text_frame.word_wrap = True
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        paragraph = text_frame.paragraphs[0]
        paragraph.alignment = align
        run = paragraph.add_run()
        run.text = text
        self._apply_run_font(run, font_size, bold, PRIMARY_FONT, TEXT_DARK)

    def _apply_run_font(self, run, font_size: float, bold: bool, font_name: str, color: RGBColor) -> None:  # 目的：统一设置文字 run 样式；功能：收敛字体、字号、粗细和颜色的重复设置；实现逻辑：所有文本最终都通过该方法落样式；可调参数：run、font_size、bold、font_name、color；默认参数及原因：调用方显式传值，避免不同分区样式互相污染。
        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color

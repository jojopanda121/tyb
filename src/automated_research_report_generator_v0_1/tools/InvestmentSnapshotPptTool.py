from __future__ import annotations

from pathlib import Path
from typing import Type

from crewai.tools import BaseTool
from pydantic import BaseModel, Field
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt
from automated_research_report_generator_v0_1.tools.investment_snapshot_ppt_a4 import (
    InvestmentSnapshotFinancialRow as _A4InvestmentSnapshotFinancialRow,
    InvestmentSnapshotOverviewProductItem as _A4InvestmentSnapshotOverviewProductItem,
    InvestmentSnapshotPptInput as _A4InvestmentSnapshotPptInput,
    InvestmentSnapshotPptTool as _A4InvestmentSnapshotPptTool,
    InvestmentSnapshotTitledItem as _A4InvestmentSnapshotTitledItem,
    ParsedFinancialTable as _A4ParsedFinancialTable,
)


"""单页投资快照 PPT 导出工具"""


FALLBACK_TEXT = "缺乏信息"
PRIMARY_FONT = "Microsoft YaHei"
TITLE_FONT = "SimHei"
SLIDE_WIDTH = Inches(13.333333)
SLIDE_HEIGHT = Inches(7.5)
BACKGROUND_COLOR = RGBColor(247, 244, 239)
BRICK_RED = RGBColor(154, 52, 28)
DEEP_RED = RGBColor(120, 38, 24)
GOLD = RGBColor(198, 149, 82)
LIGHT_BEIGE = RGBColor(239, 232, 223)
SOFT_BEIGE = RGBColor(243, 238, 232)
TEXT_DARK = RGBColor(33, 28, 24)
DASH_GOLD = RGBColor(208, 177, 130)
REQUIRED_FINANCIAL_LABELS = (
    "营业收入",
    "收入增速(%)",
    "净利润",
    "净利润增速(%)",
    "流动资产",
    "股东权益",
    "销售毛利率(%)",
    "销售净利率(%)",
)


class InvestmentSnapshotFinancialRow(BaseModel):  # 目的：定义财务行结构；功能：约束单个财务指标及其期间值；实现逻辑：使用标签加数值列表承接 agent 的结构化输出；可调参数：label、values；默认参数及原因：均为必填，避免工具在渲染期猜字段。
    label: str = Field(..., description="财务指标名称，例如营业收入")
    values: list[str] = Field(..., description="与 financial_periods 对齐的期间值列表")


class InvestmentSnapshotPptInput(BaseModel):  # 目的：定义单页 PPT 工具入参；功能：把文案和表格数据固定为可验证的结构；实现逻辑：使用明确字段替代自由文本以减少渲染歧义；可调参数：各内容字段与输出路径；默认参数及原因：全部关键字段必填，确保单页布局可以稳定生成。
    pptx_path: str = Field(..., description="输出 PPTX 文件路径")
    slide_title: str = Field(..., description="页面主标题，固定建议为 投资要点速览")
    positioning_line: str = Field(..., description="公司定位短句，会与公司名拼接为红色 banner")
    overview_points: list[str] = Field(..., description="公司概况要点，建议 3-5 条")
    financial_unit_note: str = Field(..., description="财务表单位说明，例如 亿元 或 百万元")
    financial_periods: list[str] = Field(..., description="财务期间列表，仅支持 3-4 个期间")
    financial_rows: list[InvestmentSnapshotFinancialRow] = Field(
        ...,
        description="财务数据行，使用 {label, values} 结构传入",
    )
    highlight_points: list[str] = Field(..., description="投资亮点要点，固定建议 3 条")
    risk_points: list[str] = Field(..., description="投资风险要点，固定建议 2 条")


class InvestmentSnapshotPptTool(BaseTool):  # 目的：定义单页投资快照 PPT 导出工具；功能：根据结构化输入直接绘制真实 PPTX；实现逻辑：agent 负责压缩内容，工具负责统一布局、容错和文件导出；可调参数：路径、文案、表格数据；默认参数及原因：工具名固定且返回结果作为答案，便于任务直接产出文件消息。
    name: str = "investment_snapshot_ppt_tool"
    description: str = (
        "Create a single-slide investment snapshot PowerPoint (.pptx) with a fixed "
        "investment-committee layout using only the provided structured content."
    )
    args_schema: Type[BaseModel] = InvestmentSnapshotPptInput

    def _run(  # 目的：执行 PPT 导出主流程；功能：校验输入、规范化内容并生成单页 PPTX；实现逻辑：先整理结构化数据，再统一调用绘图方法输出文件；可调参数：全部工具入参；默认参数及原因：不提供额外默认值，避免静默偏离 agent 产出。
        self,
        pptx_path: str,
        slide_title: str,
        positioning_line: str,
        overview_points: list[str],
        financial_unit_note: str,
        financial_periods: list[str],
        financial_rows: list[InvestmentSnapshotFinancialRow],
        highlight_points: list[str],
        risk_points: list[str],
    ) -> str:
        out_file = Path(pptx_path).expanduser().resolve()
        out_file.parent.mkdir(parents=True, exist_ok=True)

        company_name = self._derive_company_name_from_output_path(out_file)
        normalized_title = self._truncate_text(slide_title.strip() or "投资要点速览", max_chars=18)
        normalized_positioning = self._truncate_text(positioning_line.strip() or FALLBACK_TEXT, max_chars=44)
        normalized_overview = self._normalize_text_list(
            overview_points,
            min_items=3,
            max_items=5,
            max_chars=74,
        )
        normalized_periods = self._normalize_financial_periods(financial_periods)
        normalized_rows = self._normalize_financial_rows(
            financial_rows,
            period_count=len(normalized_periods),
        )
        normalized_highlights = self._normalize_text_list(
            highlight_points,
            min_items=3,
            max_items=3,
            max_chars=132,
        )
        normalized_risks = self._normalize_text_list(
            risk_points,
            min_items=2,
            max_items=2,
            max_chars=155,
        )

        presentation = self._build_presentation(
            company_name=company_name,
            slide_title=normalized_title,
            positioning_line=normalized_positioning,
            overview_points=normalized_overview,
            financial_unit_note=financial_unit_note.strip() or "单位未披露",
            financial_periods=normalized_periods,
            financial_rows=normalized_rows,
            highlight_points=normalized_highlights,
            risk_points=normalized_risks,
        )
        presentation.save(str(out_file))
        return f"PPT created successfully at: {out_file}"

    def _build_presentation(  # 目的：创建演示文稿对象；功能：设置画布尺寸并绘制整页内容；实现逻辑：统一创建空白页后分区调用顶部、概况、财务、亮点和风险面板；可调参数：标准化后的内容字段；默认参数及原因：单页 blank layout，减少模板依赖和主题漂移。
        self,
        company_name: str,
        slide_title: str,
        positioning_line: str,
        overview_points: list[str],
        financial_unit_note: str,
        financial_periods: list[str],
        financial_rows: list[InvestmentSnapshotFinancialRow],
        highlight_points: list[str],
        risk_points: list[str],
    ) -> Presentation:
        presentation = Presentation()
        presentation.slide_width = SLIDE_WIDTH
        presentation.slide_height = SLIDE_HEIGHT

        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = BACKGROUND_COLOR

        self._add_title_area(
            slide=slide,
            company_name=company_name,
            slide_title=slide_title,
            positioning_line=positioning_line,
        )
        self._add_overview_panel(slide=slide, overview_points=overview_points)
        self._add_financial_panel(
            slide=slide,
            financial_unit_note=financial_unit_note,
            financial_periods=financial_periods,
            financial_rows=financial_rows,
        )
        self._add_highlights_panel(slide=slide, highlight_points=highlight_points)
        self._add_risk_panel(slide=slide, risk_points=risk_points)
        return presentation

    def _derive_company_name_from_output_path(  # 目的：从固定输出文件名推导公司名；功能：避免额外增加 company_name 入参；实现逻辑：优先剥离 _investment_snapshot 后缀，失败时退回 stem；可调参数：输出路径；默认参数及原因：约定输出命名固定，适合当前任务链。
        self,
        pptx_path: Path,
    ) -> str:
        stem = pptx_path.stem.strip()
        suffix = "_investment_snapshot"
        if stem.endswith(suffix):
            candidate = stem[: -len(suffix)].strip()
            if candidate:
                return candidate
        return stem or "未知公司"

    def _normalize_financial_periods(  # 目的：规范化财务期间列表；功能：保证表头只使用 3-4 个期间且顺序稳定；实现逻辑：清洗空值后保留最新 4 个，少于 3 个则报错；可调参数：financial_periods；默认参数及原因：最少 3 期是本页财务对比的最小可读要求。
        self,
        financial_periods: list[str],
    ) -> list[str]:
        periods = [self._truncate_text(period.strip(), max_chars=18) for period in financial_periods if str(period).strip()]
        if len(periods) < 3:
            raise ValueError("financial_periods must contain at least 3 non-empty periods.")
        if len(periods) > 4:
            periods = periods[-4:]
        return periods

    def _normalize_financial_rows(  # 目的：规范化财务表行；功能：按固定指标顺序补齐或截断行数据；实现逻辑：先建立标签映射，再按 REQUIRED_FINANCIAL_LABELS 顺序输出；可调参数：financial_rows、period_count；默认参数及原因：缺失值统一补 缺乏信息，避免表格炸版或误空白。
        self,
        financial_rows: list[InvestmentSnapshotFinancialRow] | list[dict[str, object]],
        period_count: int,
    ) -> list[InvestmentSnapshotFinancialRow]:
        row_lookup: dict[str, list[str]] = {}
        for row in financial_rows:
            label, values = self._extract_financial_row_payload(row)
            if not label:
                continue
            row_lookup[label] = [
                self._truncate_text(str(value).strip() or FALLBACK_TEXT, max_chars=16)
                for value in values
            ]

        normalized_rows: list[InvestmentSnapshotFinancialRow] = []
        for label in REQUIRED_FINANCIAL_LABELS:
            values = row_lookup.get(label, [])
            if len(values) > period_count:
                values = values[-period_count:]
            if len(values) < period_count:
                values = values + [FALLBACK_TEXT] * (period_count - len(values))
            normalized_rows.append(
                InvestmentSnapshotFinancialRow(
                    label=label,
                    values=values,
                )
            )
        return normalized_rows

    def _extract_financial_row_payload(  # 目的：兼容不同来源的财务行对象；功能：同时支持 Pydantic 模型和 CrewAI 传入的 dict；实现逻辑：统一抽取 label 与 values 两个核心字段；可调参数：row；默认参数及原因：无法识别时退回空标签和空值，避免异常中断全页生成。
        self,
        row: InvestmentSnapshotFinancialRow | dict[str, object],
    ) -> tuple[str, list[object]]:
        if isinstance(row, InvestmentSnapshotFinancialRow):
            return str(row.label).strip(), list(row.values)

        if isinstance(row, dict):
            raw_values = row.get("values", [])
            values = list(raw_values) if isinstance(raw_values, list) else []
            return str(row.get("label", "")).strip(), values

        return "", []

    def _normalize_text_list(  # 目的：规范化要点列表；功能：限制条数和单条长度并在缺失时补占位；实现逻辑：先清洗空值，再按上下限截断或补足；可调参数：items、min_items、max_items、max_chars；默认参数及原因：补 缺乏信息 比留空更利于任务验收和人工检查。
        self,
        items: list[str],
        min_items: int,
        max_items: int,
        max_chars: int,
    ) -> list[str]:
        cleaned_items = [
            self._truncate_text(str(item).strip(), max_chars=max_chars)
            for item in items
            if str(item).strip()
        ]
        cleaned_items = cleaned_items[:max_items]
        while len(cleaned_items) < min_items:
            cleaned_items.append(FALLBACK_TEXT)
        return cleaned_items

    def _truncate_text(  # 目的：限制单段文字长度；功能：在不改动事实的前提下避免单页文本溢出；实现逻辑：超过阈值时尾部加省略号；可调参数：text、max_chars；默认参数及原因：长度阈值由各分区调用方决定，便于针对不同区域调优。
        self,
        text: str,
        max_chars: int,
    ) -> str:
        normalized_text = " ".join(text.split())
        if len(normalized_text) <= max_chars:
            return normalized_text
        return normalized_text[: max_chars - 1].rstrip() + "…"

    def _fit_font_size(  # 目的：按文本总长度粗调字号；功能：在固定区域内尽量保持单页不炸版；实现逻辑：根据文本总长度和最大长度做启发式缩放；可调参数：texts、base_size、min_size；默认参数及原因：优先保留基础字号，超长时逐档收缩而不是直接压到最小值。
        self,
        texts: list[str],
        base_size: int,
        min_size: int,
    ) -> int:
        total_chars = sum(len(text) for text in texts)
        longest_text = max((len(text) for text in texts), default=0)
        font_size = base_size
        if total_chars > 240 or longest_text > 90:
            font_size -= 1
        if total_chars > 320 or longest_text > 115:
            font_size -= 1
        if total_chars > 410 or longest_text > 145:
            font_size -= 1
        return max(min_size, font_size)

    def _add_title_area(  # 目的：绘制页面顶部标题区；功能：输出主标题、金色分割线和红色 banner；实现逻辑：使用文本框与色块组合出接近参考图的页首样式；可调参数：公司名、标题、定位短句；默认参数及原因：标题区高度固定，保证下方四区布局稳定。
        self,
        slide,
        company_name: str,
        slide_title: str,
        positioning_line: str,
    ) -> None:
        self._add_textbox(
            slide=slide,
            left=Inches(0.55),
            top=Inches(0.18),
            width=Inches(4.2),
            height=Inches(0.52),
            text=slide_title,
            font_size=28,
            bold=True,
            font_name=TITLE_FONT,
            color=TEXT_DARK,
            align=PP_ALIGN.LEFT,
            margin_left=2,
            margin_right=2,
            margin_top=0,
            margin_bottom=0,
        )

        line_shape = slide.shapes.add_textbox(
            Inches(0.3),
            Inches(0.84),
            Inches(12.9),
            Inches(0.08),
        )
        line_shape.fill.solid()
        line_shape.fill.fore_color.rgb = GOLD
        line_shape.line.fill.background()

        triangle = slide.shapes.add_textbox(
            Inches(13.02),
            Inches(0.84),
            Inches(0.22),
            Inches(0.08),
        )
        triangle.fill.solid()
        triangle.fill.fore_color.rgb = GOLD
        triangle.line.fill.background()
        triangle.rotation = 35

        banner_text = f"{company_name}：{positioning_line}"
        self._add_textbox(
            slide=slide,
            left=Inches(0.12),
            top=Inches(1.18),
            width=Inches(12.9),
            height=Inches(0.48),
            text=banner_text,
            font_size=20,
            bold=True,
            font_name=TITLE_FONT,
            color=RGBColor(255, 255, 255),
            fill_color=BRICK_RED,
            line_color=BRICK_RED,
            align=PP_ALIGN.LEFT,
            margin_left=12,
            margin_right=12,
            margin_top=2,
            margin_bottom=0,
        )

    def _add_overview_panel(  # 目的：绘制公司概况面板；功能：承载公司和产品简介的 3-5 条要点；实现逻辑：先画虚线框，再叠加标题标签和多段正文文本框；可调参数：overview_points；默认参数及原因：面板尺寸固定，对应参考图上半区左侧区域。
        self,
        slide,
        overview_points: list[str],
    ) -> None:
        panel_left = Inches(0.12)
        panel_top = Inches(1.94)
        panel_width = Inches(6.34)
        panel_height = Inches(2.48)

        self._add_panel_frame(
            slide=slide,
            left=panel_left,
            top=panel_top,
            width=panel_width,
            height=panel_height,
        )
        self._add_section_label(
            slide=slide,
            text="公司概况",
            left=Inches(2.25),
            top=Inches(1.72),
            width=Inches(1.95),
        )

        body_font_size = self._fit_font_size(overview_points, base_size=11, min_size=9)
        self._add_list_textbox(
            slide=slide,
            left=panel_left + Inches(0.1),
            top=panel_top + Inches(0.18),
            width=panel_width - Inches(0.2),
            height=panel_height - Inches(0.26),
            lines=overview_points,
            font_size=body_font_size,
            bullet_prefix="•",
            color=TEXT_DARK,
        )

    def _add_financial_panel(  # 目的：绘制财务数据面板；功能：呈现期间表头和 8 行核心指标；实现逻辑：标题标签、表头文本框和一个 8 行表格分层组合；可调参数：单位、期间、财务行；默认参数及原因：表头独立于表格，便于满足 8 行数据体的验收要求。
        self,
        slide,
        financial_unit_note: str,
        financial_periods: list[str],
        financial_rows: list[InvestmentSnapshotFinancialRow],
    ) -> None:
        panel_left = Inches(6.56)
        panel_top = Inches(1.94)
        panel_width = Inches(6.46)
        panel_height = Inches(2.48)
        period_count = len(financial_periods)

        self._add_panel_frame(
            slide=slide,
            left=panel_left,
            top=panel_top,
            width=panel_width,
            height=panel_height,
            show_fill=False,
        )
        self._add_section_label(
            slide=slide,
            text="财务数据",
            left=Inches(8.92),
            top=Inches(1.72),
            width=Inches(1.95),
        )

        content_left = panel_left + Inches(0.06)
        content_width = panel_width - Inches(0.12)
        first_column_width = Inches(1.7)
        other_column_width = int((content_width - first_column_width) / period_count)
        column_widths = [first_column_width] + [other_column_width] * period_count

        header_top = panel_top + Inches(0.16)
        header_height = Inches(0.34)
        header_labels = [f"项目（{financial_unit_note}）"] + financial_periods

        running_left = content_left
        for index, header_text in enumerate(header_labels):
            self._add_textbox(
                slide=slide,
                left=running_left,
                top=header_top,
                width=column_widths[index],
                height=header_height,
                text=header_text,
                font_size=10,
                bold=True,
                font_name=PRIMARY_FONT,
                color=TEXT_DARK,
                align=PP_ALIGN.CENTER if index > 0 else PP_ALIGN.LEFT,
                margin_left=4 if index == 0 else 2,
                margin_right=2,
                margin_top=1,
                margin_bottom=0,
            )
            running_left += column_widths[index]

        top_rule = slide.shapes.add_textbox(
            content_left,
            header_top - Inches(0.02),
            content_width,
            Inches(0.02),
        )
        top_rule.fill.solid()
        top_rule.fill.fore_color.rgb = BRICK_RED
        top_rule.line.fill.background()

        bottom_rule = slide.shapes.add_textbox(
            content_left,
            header_top + header_height + Inches(0.01),
            content_width,
            Inches(0.02),
        )
        bottom_rule.fill.solid()
        bottom_rule.fill.fore_color.rgb = BRICK_RED
        bottom_rule.line.fill.background()

        table_top = panel_top + Inches(0.54)
        table_height = panel_height - Inches(0.62)
        table_shape = slide.shapes.add_table(
            rows=len(financial_rows),
            cols=1 + period_count,
            left=content_left,
            top=table_top,
            width=content_width,
            height=table_height,
        )
        table = table_shape.table
        table.first_row = False
        table.horz_banding = False
        table.vert_banding = False

        for column_index, column_width in enumerate(column_widths):
            table.columns[column_index].width = column_width

        row_height = int(table_height / max(1, len(financial_rows)))
        for row_index, row_payload in enumerate(financial_rows):
            table.rows[row_index].height = row_height
            self._set_table_cell_text(
                cell=table.cell(row_index, 0),
                text=row_payload.label,
                font_size=10,
                bold=True,
                align=PP_ALIGN.LEFT,
            )
            for value_index, value in enumerate(row_payload.values, start=1):
                self._set_table_cell_text(
                    cell=table.cell(row_index, value_index),
                    text=value,
                    font_size=10,
                    bold=False,
                    align=PP_ALIGN.CENTER,
                )

    def _add_highlights_panel(  # 目的：绘制投资亮点面板；功能：承载 3 条投资亮点；实现逻辑：使用整宽虚线框和多段要点文本提升信息密度；可调参数：highlight_points；默认参数及原因：固定三条，便于基金内部快速浏览和比较。
        self,
        slide,
        highlight_points: list[str],
    ) -> None:
        panel_left = Inches(0.12)
        panel_top = Inches(4.56)
        panel_width = Inches(12.91)
        panel_height = Inches(1.62)

        self._add_panel_frame(
            slide=slide,
            left=panel_left,
            top=panel_top,
            width=panel_width,
            height=panel_height,
        )
        self._add_section_label(
            slide=slide,
            text="投资亮点",
            left=Inches(4.7),
            top=Inches(4.34),
            width=Inches(2.4),
        )

        body_font_size = self._fit_font_size(highlight_points, base_size=11, min_size=9)
        self._add_list_textbox(
            slide=slide,
            left=panel_left + Inches(0.12),
            top=panel_top + Inches(0.16),
            width=panel_width - Inches(0.24),
            height=panel_height - Inches(0.22),
            lines=highlight_points,
            font_size=body_font_size,
            bullet_prefix="•",
            color=TEXT_DARK,
        )

    def _add_risk_panel(  # 目的：绘制底部风险条带；功能：突出 2 条最关键投资风险；实现逻辑：左侧标签块加右侧正文条带，形成底部总结区域；可调参数：risk_points；默认参数及原因：固定两条，避免风险信息稀释和阅读拥挤。
        self,
        slide,
        risk_points: list[str],
    ) -> None:
        panel_left = Inches(0.12)
        panel_top = Inches(6.34)
        panel_width = Inches(12.91)
        panel_height = Inches(1.03)
        label_width = Inches(1.95)

        self._add_textbox(
            slide=slide,
            left=panel_left,
            top=panel_top,
            width=label_width,
            height=panel_height,
            text="投资风险",
            font_size=18,
            bold=True,
            font_name=TITLE_FONT,
            color=BRICK_RED,
            fill_color=LIGHT_BEIGE,
            line_color=LIGHT_BEIGE,
            align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
            margin_left=2,
            margin_right=2,
            margin_top=0,
            margin_bottom=0,
        )
        self._add_textbox(
            slide=slide,
            left=panel_left + label_width,
            top=panel_top,
            width=panel_width - label_width,
            height=panel_height,
            text="",
            font_size=10,
            bold=False,
            fill_color=SOFT_BEIGE,
            line_color=SOFT_BEIGE,
        )

        body_font_size = self._fit_font_size(risk_points, base_size=11, min_size=9)
        self._add_list_textbox(
            slide=slide,
            left=panel_left + label_width + Inches(0.08),
            top=panel_top + Inches(0.08),
            width=panel_width - label_width - Inches(0.16),
            height=panel_height - Inches(0.16),
            lines=risk_points,
            font_size=body_font_size,
            bullet_prefix="•",
            color=TEXT_DARK,
            space_after=1,
        )

    def _add_panel_frame(  # 目的：绘制虚线内容框；功能：统一概况、财务和亮点区的边框风格；实现逻辑：使用无圆角文本框作为容器边框；可调参数：位置尺寸和是否填充；默认参数及原因：默认浅底色加金色虚线，接近参考页的基金投决卡片风格。
        self,
        slide,
        left,
        top,
        width,
        height,
        show_fill: bool = True,
    ) -> None:
        frame = slide.shapes.add_textbox(left, top, width, height)
        if show_fill:
            frame.fill.solid()
            frame.fill.fore_color.rgb = RGBColor(250, 248, 244)
        else:
            frame.fill.background()
        frame.line.color.rgb = DASH_GOLD
        frame.line.width = Pt(1.2)
        frame.line.dash_style = MSO_LINE_DASH_STYLE.DASH

    def _add_section_label(  # 目的：绘制分区标签；功能：为每个内容面板提供居中的标题签；实现逻辑：使用浅米色文本框压在边框上方；可调参数：文字和位置；默认参数及原因：统一 14pt 粗体，增强单页层次感。
        self,
        slide,
        text: str,
        left,
        top,
        width,
    ) -> None:
        self._add_textbox(
            slide=slide,
            left=left,
            top=top,
            width=width,
            height=Inches(0.32),
            text=text,
            font_size=14,
            bold=True,
            font_name=TITLE_FONT,
            color=BRICK_RED,
            fill_color=LIGHT_BEIGE,
            line_color=LIGHT_BEIGE,
            align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
            margin_left=0,
            margin_right=0,
            margin_top=0,
            margin_bottom=0,
        )

    def _add_list_textbox(  # 目的：绘制多段要点文本框；功能：统一 bullet 风格、换行和段间距；实现逻辑：逐条建立段落而不是拼接长字符串，便于后续微调；可调参数：位置、要点、字号和 bullet 前缀；默认参数及原因：默认使用圆点前缀，贴近研究纪要阅读习惯。
        self,
        slide,
        left,
        top,
        width,
        height,
        lines: list[str],
        font_size: int,
        bullet_prefix: str,
        color: RGBColor,
        space_after: int = 2,
    ) -> None:
        textbox = slide.shapes.add_textbox(left, top, width, height)
        textbox.fill.background()
        textbox.line.fill.background()

        text_frame = textbox.text_frame
        text_frame.clear()
        text_frame.word_wrap = True
        text_frame.vertical_anchor = MSO_ANCHOR.TOP
        text_frame.margin_left = Pt(0)
        text_frame.margin_right = Pt(0)
        text_frame.margin_top = Pt(0)
        text_frame.margin_bottom = Pt(0)

        for index, line in enumerate(lines):
            paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
            paragraph.alignment = PP_ALIGN.LEFT
            paragraph.space_after = Pt(space_after)
            run = paragraph.add_run()
            run.text = f"{bullet_prefix} {line}"
            self._apply_run_font(
                run=run,
                font_size=font_size,
                bold=False,
                font_name=PRIMARY_FONT,
                color=color,
            )

    def _add_textbox(  # 目的：封装单段文本框创建；功能：统一文本框的字体、填充、边框和内边距设置；实现逻辑：把常用文本框参数收敛到一个方法里减少重复；可调参数：位置、样式、对齐和内边距；默认参数及原因：默认无底色无边框，适合大多数正文场景。
        self,
        slide,
        left,
        top,
        width,
        height,
        text: str,
        font_size: int,
        bold: bool,
        font_name: str = PRIMARY_FONT,
        color: RGBColor = TEXT_DARK,
        fill_color: RGBColor | None = None,
        line_color: RGBColor | None = None,
        align: PP_ALIGN = PP_ALIGN.LEFT,
        anchor: MSO_ANCHOR = MSO_ANCHOR.TOP,
        margin_left: int = 6,
        margin_right: int = 6,
        margin_top: int = 4,
        margin_bottom: int = 4,
    ):
        textbox = slide.shapes.add_textbox(left, top, width, height)
        if fill_color is None:
            textbox.fill.background()
        else:
            textbox.fill.solid()
            textbox.fill.fore_color.rgb = fill_color

        if line_color is None:
            textbox.line.fill.background()
        else:
            textbox.line.color.rgb = line_color
            textbox.line.width = Pt(0.8)

        text_frame = textbox.text_frame
        text_frame.clear()
        text_frame.word_wrap = True
        text_frame.vertical_anchor = anchor
        text_frame.margin_left = Pt(margin_left)
        text_frame.margin_right = Pt(margin_right)
        text_frame.margin_top = Pt(margin_top)
        text_frame.margin_bottom = Pt(margin_bottom)

        paragraph = text_frame.paragraphs[0]
        paragraph.alignment = align
        run = paragraph.add_run()
        run.text = text
        self._apply_run_font(
            run=run,
            font_size=font_size,
            bold=bold,
            font_name=font_name,
            color=color,
        )
        return textbox

    def _set_table_cell_text(  # 目的：统一设置表格单元格文本；功能：控制字体、对齐、填充和边距，降低默认主题干扰；实现逻辑：逐格写入单个 run，避免遗留模板样式；可调参数：cell、text、font_size、bold、align；默认参数及原因：默认白底深色字，更适合财务表阅读。
        self,
        cell,
        text: str,
        font_size: int,
        bold: bool,
        align: PP_ALIGN,
    ) -> None:
        cell.text = ""
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = Pt(3)
        cell.margin_right = Pt(3)
        cell.margin_top = Pt(1)
        cell.margin_bottom = Pt(1)

        text_frame = cell.text_frame
        text_frame.word_wrap = True
        paragraph = text_frame.paragraphs[0]
        paragraph.alignment = align
        run = paragraph.add_run()
        run.text = text
        self._apply_run_font(
            run=run,
            font_size=font_size,
            bold=bold,
            font_name=PRIMARY_FONT,
            color=TEXT_DARK,
        )

    def _apply_run_font(  # 目的：统一设置文字 run 样式；功能：收敛字体、字号、粗细和颜色的重复设置；实现逻辑：所有文本最终都通过该方法落样式；可调参数：run、font_size、bold、font_name、color；默认参数及原因：调用方显式传值，避免不同分区样式互相污染。
        self,
        run,
        font_size: int,
        bold: bool,
        font_name: str,
        color: RGBColor,
    ) -> None:
        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color


InvestmentSnapshotFinancialRow = _A4InvestmentSnapshotFinancialRow
InvestmentSnapshotOverviewProductItem = _A4InvestmentSnapshotOverviewProductItem
InvestmentSnapshotTitledItem = _A4InvestmentSnapshotTitledItem
InvestmentSnapshotPptInput = _A4InvestmentSnapshotPptInput
ParsedFinancialTable = _A4ParsedFinancialTable
InvestmentSnapshotPptTool = _A4InvestmentSnapshotPptTool

from __future__ import annotations

import json
from pathlib import Path
from time import time_ns

from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from automated_research_report_generator_v0_1.tools.investment_snapshot_ppt_a4 import (
    FALLBACK_TEXT,
    InvestmentSnapshotOverviewProductItem,
    InvestmentSnapshotPptTool,
    InvestmentSnapshotTitledItem,
    TABLE_ROW_FILL_PRIMARY,
    TABLE_ROW_FILL_SECONDARY,
)


"""横向 A4 单页投资快照 PPT 工具测试脚本"""


TEST_OUTPUT_DIR = Path.cwd() / ".cache" / "investment_snapshot_ppt_test"
TEST_COMPANY_NAME = "测试公司股份有限公司"
EXPECTED_SLIDE_WIDTH = Inches(11.69)
EXPECTED_SLIDE_HEIGHT = Inches(8.27)
ALIAS_MARKDOWN = """# 测试公司股份有限公司 财务分析
## 1. 核心财务数据总表
以下数据基于公司披露口径整理，币种为人民币，单位为百万元。
| 指标 | 2023年12月31日 | 2024年12月31日 | 2025年12月31日 | 数据来源（报表/附注/页码） | 备注 |
| :--- | :--- | :--- | :--- | :--- | :--- |
| 营业收入 | 100 | 150 | 210 | 利润表 / P10 | 单位：百万元 |
| 收入增速(%) | 缺乏信息 | 50.0% | 40.0% | 推导 | |
| 毛利润 | 35 | 55 | 79.8 | 利润表 / P10 | |
| 销售毛利率(%) | 35.0% | 36.7% | 38.0% | 附注 / P18 | |
| 净利润 | 10 | 18 | 27 | 利润表 / P10 | |
| 销售净利率(%) | 10.0% | 12.0% | 12.9% | 推导 | |
| 净利润增速(%) | 缺乏信息 | 80.0% | 50.0% | 推导 | |
| 资产总计 | 400 | 560 | 680 | 资产负债表 / P12 | |
| 资产负债率 | 52.5% | 58.9% | 61.8% | 资产负债表 / P12 | |
| ROE | 9.0% | 12.5% | 14.2% | 财务比率 / P20 | |
## 2. 财务趋势分析
- 略
"""
DERIVED_MARKDOWN = """# 测试公司股份有限公司 财务分析
## 1. 核心财务数据总表
| 指标 | 2023年12月31日 | 2024年12月31日 | 2025年12月31日 | 数据来源（报表/附注/页码） | 备注 |
| :--- | :--- | :--- | :--- | :--- | :--- |
| 营业收入 | 100 | 150 | 210 | 利润表 / P10 | 单位：人民币百万元 |
| 毛利率 | 35.0% | 36.7% | 38.0% | 附注 / P18 | |
| 净利润 | -10 | 15 | 21 | 利润表 / P10 | |
| 总资产 | 400 | 560 | 700 | 资产负债表 / P12 | |
| 股东权益 | 190 | 230 | 260 | 资产负债表 / P12 | |
## 2. 财务趋势分析
- 略
"""
LONG_FINANCIAL_MARKDOWN = """# 测试公司股份有限公司 财务分析
## 1. 核心财务数据总表
以下数据基于公司披露口径整理，币种：人民币；单位：元。
| 指标 | 2024年9月30日 | 2024年12月31日 | 2025年9月30日 | 数据来源（报表/附注/页码） | 备注 |
| :--- | :--- | :--- | :--- | :--- | :--- |
| 营业收入 | 123,456,789,012.3456 | 234,567,890,123.4567 | 345,678,901,234.5678 | 利润表 / P10 | 单位：人民币元 |
| 毛利率 | 35.01234% | 36.67891% | 38.00123% | 附注 / P18 | |
| 净利润 | (10,987,654,321.1234) | 15,876,543,210.2345 | 21,765,432,109.3456 | 利润表 / P10 | |
| 总资产 | 400,123,456,789.1234 | 560,234,567,890.2345 | 700,345,678,901.3456 | 资产负债表 / P12 | |
| 股东权益 | 190,123,456,789.1111 | 230,234,567,890.2222 | 260,345,678,901.3333 | 资产负债表 / P12 | |
## 2. 财务趋势分析
- 略
"""


def build_overview_product_items() -> list[InvestmentSnapshotOverviewProductItem]:  # 目的：构造公司概况产品样例；功能：为概况区提供固定 3 条“产品名 + 描述”；实现逻辑：使用确定文本便于后续断言粗体 run 和文本命中；可调参数：无；默认参数及原因：固定样例最有利于重复运行和排查版式问题。
    return [
        InvestmentSnapshotOverviewProductItem(name="LightSphere X", description="面向 AI 超节点集群的高带宽低时延光互连交换设备，用于缓解机内与机间通信瓶颈。"),
        InvestmentSnapshotOverviewProductItem(name="PACE 3", description="用于生成式 AI 训练与推理的光电混合算力加速卡与模组，强调能效和系统级吞吐提升。"),
        InvestmentSnapshotOverviewProductItem(name="Gazelle", description="帮助客户完成评估验证、软件迁移和系统联调的开发套件，缩短从 PoC 到规模部署的周期。"),
    ]


def build_highlight_items() -> list[InvestmentSnapshotTitledItem]:  # 目的：构造投资亮点样例；功能：为亮点区提供 3 条“短标题 + 说明句”内容；实现逻辑：使用固定结构化输入验证富文本而非 bullet 输出；可调参数：无；默认参数及原因：固定值便于验证条数上限和粗体标题 run。
    return [
        InvestmentSnapshotTitledItem(title="技术壁垒", detail="公司以光互连与光计算协同架构切入 AI 基础设施瓶颈，核心器件、自研系统设计和客户联调能力共同构成较高导入门槛，使竞争不再停留在单一器件层面。"),
        InvestmentSnapshotTitledItem(title="商业放量", detail="随着头部客户从测试验证转向批量部署，收入结构正从项目制研发服务逐步转向标准化硬件和系统销售，单客户价值量、复购概率和交付可预测性同步提升。"),
        InvestmentSnapshotTitledItem(title="生态绑定", detail="公司与多家 GPU、服务器和数据中心基础设施厂商开展联合设计导入，验证链条长但一旦进入正式采购名单，后续替换成本高、合作黏性强。"),
    ]


def build_risk_items() -> list[InvestmentSnapshotTitledItem]:  # 目的：构造投资风险样例；功能：为风险区提供 2 条“短标题 + 说明句”内容；实现逻辑：使用固定结构验证风险条数和底部条带版式；可调参数：无；默认参数及原因：控制在 2 条以内以贴合目标页面密度。
    return [
        InvestmentSnapshotTitledItem(title="量产落地风险", detail="硬件产品仍处于放量早期，若良率爬坡、供应链协同或客户机房导入节奏不及预期，订单确认和收入兑现都可能出现阶段性延后。"),
        InvestmentSnapshotTitledItem(title="资金消耗压力", detail="研发投入、测试验证和客户交付前置投入较高，若资本市场融资窗口收紧或新增订单回款节奏变慢，短期现金流压力可能被放大。"),
    ]


def collect_slide_texts(slide) -> list[str]:  # 目的：提取页面全部文本；功能：统一收集文本框和表格单元格文字用于断言；实现逻辑：遍历 shape 中的 text_frame 和 table 两类对象；可调参数：slide；默认参数及原因：返回扁平文本列表，便于简单包含判断。
    texts: list[str] = []
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False) and shape.text_frame is not None:
            text = shape.text_frame.text.strip()
            if text:
                texts.append(text)
        if getattr(shape, "has_table", False):
            for row in shape.table.rows:
                for cell in row.cells:
                    text = cell.text.strip()
                    if text:
                        texts.append(text)
    return texts


def collect_bold_run_texts(slide) -> list[str]:  # 目的：提取页面粗体 run 文本；功能：验证产品名、亮点标题和风险标题被单独设置为粗体；实现逻辑：遍历所有文本框段落中的 runs 并筛选 bold=True；可调参数：slide；默认参数及原因：只返回非空文本，便于直接做包含断言。
    bold_texts: list[str] = []
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False) or shape.text_frame is None:
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if run.text.strip() and bool(run.font.bold):
                    bold_texts.append(run.text.strip())
    return bold_texts


def rows_to_lookup(rows) -> dict[str, list[str]]:  # 目的：把财务行列表转成查找表；功能：便于测试中直接按标签断言计算结果；实现逻辑：遍历财务行对象并按 label 建立 values 映射；可调参数：rows；默认参数及原因：返回字典最利于做多指标精确比较。
    return {row.label: list(row.values) for row in rows}


def assert_alias_mapping() -> None:  # 目的：验证旧口径或别名口径的财务行能被正确映射；功能：覆盖营业收入增长率、毛利率、净利率、净利润增长率、总资产、资产负债率和 ROE 的别名逻辑；实现逻辑：用别名样例 Markdown 调用工具内部解析和映射；可调参数：无；默认参数及原因：直接断言核心映射结果，定位问题最快。
    tool = InvestmentSnapshotPptTool()
    parsed = tool._parse_financial_markdown(ALIAS_MARKDOWN)
    rows = rows_to_lookup(tool._build_snapshot_financial_rows(parsed))

    assert parsed.display_periods == ["2023/12/31", "2024/12/31", "2025/12/31"], "期间应统一转成 YYYY/MM/DD。"
    assert parsed.unit_note == "人民币，百万元", "表头应同时展示币种和单位。"
    assert rows["营业收入"] == ["100", "150", "210"], "营业收入应直接命中。"
    assert rows["营业收入增长率"] == [FALLBACK_TEXT, "50.0%", "40.0%"], "营业收入增长率应接受旧别名映射。"
    assert rows["毛利率"] == ["35.0%", "36.7%", "38.0%"], "毛利率应接受销售毛利率别名。"
    assert rows["净利率"] == ["10.0%", "12.0%", "12.9%"], "净利率应接受销售净利率别名。"
    assert rows["净利润增长率"] == [FALLBACK_TEXT, "80.0%", "50.0%"], "净利润增长率应接受旧别名映射。"
    assert rows["总资产"] == ["400", "560", "680"], "总资产应接受资产总计别名。"
    assert rows["资产负债率"] == ["52.5%", "58.9%", "61.8%"], "资产负债率应接受直接映射。"
    assert rows["ROE"] == ["9.0%", "12.5%", "14.2%"], "ROE 应保留显式披露值。"


def assert_derived_metrics() -> None:  # 目的：验证缺失指标的有限代码推导逻辑；功能：覆盖营业收入增长率、净利率、净利润增长率和资产负债率的推导，以及 ROE 缺失回填；实现逻辑：用缺少显式比率行的 Markdown 样例走完整解析链；可调参数：无；默认参数及原因：直接围绕新增推导规则做硬断言，最能体现本次改动价值。
    tool = InvestmentSnapshotPptTool()
    parsed = tool._parse_financial_markdown(DERIVED_MARKDOWN)
    rows = rows_to_lookup(tool._build_snapshot_financial_rows(parsed))

    assert parsed.unit_note == "人民币，百万元", "金额单位写成“人民币百万元”时应自动拆出币种和单位。"
    assert rows["营业收入增长率"] == [FALLBACK_TEXT, "50.0%", "40.0%"], "营业收入增长率应可由营业收入推导。"
    assert rows["净利率"] == ["-10.0%", "10.0%", "10.0%"], "净利率应可由净利润与营业收入推导。"
    assert rows["净利润增长率"] == [FALLBACK_TEXT, "250.0%", "40.0%"], "净利润增长率应可由净利润推导。"
    assert rows["资产负债率"] == ["52.5%", "58.9%", "62.9%"], "资产负债率应可由总资产和股东权益推导。"
    assert rows["ROE"] == [FALLBACK_TEXT, FALLBACK_TEXT, FALLBACK_TEXT], "缺失 ROE 时应明确回填缺乏信息。"


def assert_financial_values_not_truncated() -> None:  # 目的：验证财务表值不会因版式原因被工具裁剪；功能：覆盖超长数值和非常规期间标签的保真输出；实现逻辑：用长数值 Markdown 样例走解析和映射链，再断言完整字符串被保留；可调参数：无；默认参数及原因：财务值精度是刚性要求，必须用回归测试锁死。
    tool = InvestmentSnapshotPptTool()
    parsed = tool._parse_financial_markdown(LONG_FINANCIAL_MARKDOWN)
    rows = rows_to_lookup(tool._build_snapshot_financial_rows(parsed))

    assert parsed.display_periods == ["2024/09/30", "2024/12/31", "2025/09/30"], "非常规期间标签应尽量标准化为 YYYY/MM/DD。"
    assert parsed.unit_note == "人民币，元", "显式分开的币种和单位应被组合展示。"
    assert rows["营业收入"] == ["123,456,789,012.3456", "234,567,890,123.4567", "345,678,901,234.5678"], "超长营业收入值应完整保留。"
    assert rows["毛利率"] == ["35.01234%", "36.67891%", "38.00123%"], "超长百分比值应完整保留。"
    assert rows["总资产"] == ["400,123,456,789.1234", "560,234,567,890.2345", "700,345,678,901.3456"], "超长资产值应完整保留。"


def assert_meta_note_from_llm_style_sentence() -> None:  # 目的：验证 snapshot 工具能兼容 LLM 常见的自由句式元数据描述；功能：覆盖“币种为人民币，单位为千元（除百分比外）”这种非冒号口径；实现逻辑：直接喂入与真实财务分析文件一致的句式并断言组合表头口径；可调参数：无；默认参数及原因：这是当前真实故障场景，必须单独锁成回归测试。
    tool = InvestmentSnapshotPptTool()
    markdown = """# 示例公司 财务分析
## 1. 核心财务数据总表
以下数据基于公司招股书披露的综合财务报表，币种为人民币，单位为千元（除百分比外）。
| 指标 | 2023年12月31日 | 2024年12月31日 | 2025年12月31日 | 数据来源（页码） | 备注 |
| :--- | :--- | :--- | :--- | :--- | :--- |
| 营业收入 | 1 | 2 | 3 | P10 | - |
| 毛利率 | 10.0% | 20.0% | 30.0% | P12 | - |
| 净利润 | 0.1 | 0.2 | 0.3 | P10 | - |
| 总资产 | 10 | 20 | 30 | P11 | - |
| 股东权益 | 6 | 8 | 10 | P11 | - |
## 2. 财务趋势分析
- 略
"""
    parsed = tool._parse_financial_markdown(markdown)
    assert parsed.unit_note == "人民币，千元", "自由句式里的币种和单位应被正确抽取并组合。"


def assert_ppt_structure(pptx_path: Path) -> dict[str, object]:  # 目的：验证生成 PPT 的结构；功能：检查页数、尺寸、关键标题、表格、文本区块和粗体 run 是否存在；实现逻辑：重新打开 PPT 后做硬断言并返回摘要结果；可调参数：pptx_path；默认参数及原因：直接抛出 AssertionError，便于命令行快速失败。
    presentation = Presentation(str(pptx_path))
    assert len(presentation.slides) == 1, "生成的 PPT 应只包含 1 张页面。"
    assert presentation.slide_width == EXPECTED_SLIDE_WIDTH, "PPT 宽度应为横向 A4。"
    assert presentation.slide_height == EXPECTED_SLIDE_HEIGHT, "PPT 高度应为横向 A4。"

    slide = presentation.slides[0]
    texts = collect_slide_texts(slide)
    merged_text = "\n".join(texts)
    bold_runs = collect_bold_run_texts(slide)

    assert f"投资要点速览——{TEST_COMPANY_NAME}" in merged_text, "页面主标题应包含固定标题和公司名。"
    assert "光电混合算力平台先行者，面向 AI 基础设施放量" in merged_text, "红色 banner 中应仅包含公司定位短句。"
    assert "公司概况" in merged_text, "公司概况分区标题缺失。"
    assert "财务数据" in merged_text, "财务数据分区标题缺失。"
    assert "投资亮点" in merged_text, "投资亮点分区标题缺失。"
    assert "投资风险" in merged_text, "投资风险分区标题缺失。"
    assert "光电混合算力平台" in merged_text, "公司摘要正文缺失。"
    assert "技术壁垒" in merged_text, "投资亮点正文缺失。"
    assert "量产落地风险" in merged_text, "投资风险正文缺失。"

    tables = [shape.table for shape in slide.shapes if getattr(shape, "has_table", False)]
    assert len(tables) == 1, "页面内应存在且仅存在 1 个财务表格。"
    assert len(tables[0].rows) == 8, "财务表格应包含 8 行核心指标。"
    assert len(tables[0].columns) == 4, "财务表格应包含 1 列指标名和 3 列期间值。"
    assert tables[0].cell(0, 0).text_frame.paragraphs[0].alignment == PP_ALIGN.LEFT, "财务表正文首列应左对齐。"
    assert tables[0].cell(0, 1).text_frame.paragraphs[0].alignment == PP_ALIGN.RIGHT, "财务表正文日期列应右对齐。"
    assert tables[0].cell(0, 0).vertical_anchor == MSO_ANCHOR.MIDDLE, "财务表正文首列应上下居中。"
    assert tables[0].cell(0, 1).vertical_anchor == MSO_ANCHOR.MIDDLE, "财务表正文日期列应上下居中。"
    header_text_lookup = {
        shape.text_frame.text.strip(): shape
        for shape in slide.shapes
        if getattr(shape, "has_text_frame", False) and shape.text_frame is not None and shape.text_frame.text.strip()
    }
    body_value_font_size = tables[0].cell(0, 1).text_frame.paragraphs[0].runs[0].font.size
    for header_text in ("项目（人民币，百万元）", "2023/12/31", "2024/12/31", "2025/12/31"):
        assert header_text in header_text_lookup, f"财务表头 {header_text} 缺失。"
        header_shape = header_text_lookup[header_text]
        assert header_shape.text_frame.paragraphs[0].alignment == PP_ALIGN.CENTER, f"财务表头 {header_text} 应居中对齐。"
        assert header_shape.text_frame.vertical_anchor == MSO_ANCHOR.MIDDLE, f"财务表头 {header_text} 应上下居中。"
    for date_header_text in ("2023/12/31", "2024/12/31", "2025/12/31"):
        date_header_shape = header_text_lookup[date_header_text]
        assert date_header_shape.text_frame.paragraphs[0].runs[0].font.size == body_value_font_size, f"日期表头 {date_header_text} 字号应与财务正文一致。"
    assert header_text_lookup[f"投资要点速览——{TEST_COMPANY_NAME}"].text_frame.vertical_anchor == MSO_ANCHOR.MIDDLE, "主标题应上下居中。"
    assert header_text_lookup["光电混合算力平台先行者，面向 AI 基础设施放量"].text_frame.vertical_anchor == MSO_ANCHOR.MIDDLE, "banner 文案应上下居中。"
    summary_shape = next(shape for shape in slide.shapes if getattr(shape, "has_text_frame", False) and "公司聚焦光电混合算力平台" in shape.text_frame.text)
    highlights_shape = next(shape for shape in slide.shapes if getattr(shape, "has_text_frame", False) and "技术壁垒：" in shape.text_frame.text)
    risk_shape = next(shape for shape in slide.shapes if getattr(shape, "has_text_frame", False) and "量产落地风险：" in shape.text_frame.text)
    assert summary_shape.text_frame.vertical_anchor == MSO_ANCHOR.MIDDLE, "公司概况正文应上下居中。"
    assert highlights_shape.text_frame.vertical_anchor == MSO_ANCHOR.MIDDLE, "投资亮点正文应上下居中。"
    assert risk_shape.text_frame.vertical_anchor == MSO_ANCHOR.MIDDLE, "投资风险正文应上下居中。"
    assert tables[0].cell(0, 0).fill.fore_color.rgb == TABLE_ROW_FILL_PRIMARY, "财务表首行应使用浅色底纹。"
    assert tables[0].cell(1, 0).fill.fore_color.rgb == TABLE_ROW_FILL_SECONDARY, "财务表第二行应使用交替底纹。"
    assert any(text == "LightSphere X：" for text in bold_runs), "产品名称应以独立粗体 run 呈现。"
    assert any(text == "技术壁垒：" for text in bold_runs), "投资亮点标题应以独立粗体 run 呈现。"
    assert any(text == "量产落地风险：" for text in bold_runs), "投资风险标题应以独立粗体 run 呈现。"

    return {
        "slide_count": len(presentation.slides),
        "slide_width": presentation.slide_width,
        "slide_height": presentation.slide_height,
        "table_count": len(tables),
        "financial_row_count": len(tables[0].rows),
        "bold_run_samples": [text for text in bold_runs if text in {"LightSphere X：", "技术壁垒：", "量产落地风险："}],
    }


def run_tool_smoke_test() -> dict[str, object]:  # 目的：执行工具级冒烟测试；功能：生成样例 PPT 并复开校验其结构正确；实现逻辑：准备固定输入、调用工具、再做二次读取断言；可调参数：无；默认参数及原因：使用仓库内 `.cache` 目录输出，避免污染正式产物目录。
    run_output_dir = TEST_OUTPUT_DIR / f"run_{time_ns()}"
    run_output_dir.mkdir(parents=True, exist_ok=True)
    pptx_path = run_output_dir / f"{TEST_COMPANY_NAME}_investment_snapshot.pptx"

    tool = InvestmentSnapshotPptTool()
    result_message = tool._run(
        pptx_path=str(pptx_path),
        slide_title="投资要点速览",
        positioning_line="光电混合算力平台先行者，面向 AI 基础设施放量",
        overview_summary="公司聚焦光电混合算力平台，通过光互连与光计算产品解决 AI 基础设施在带宽、时延和能效上的核心瓶颈，既具备底层器件与系统协同能力，也正进入客户验证向规模采购切换的关键阶段，兼具技术稀缺性与商业放量吸引力。",
        overview_product_items=build_overview_product_items(),
        financial_source_markdown=DERIVED_MARKDOWN,
        highlight_items=build_highlight_items(),
        risk_items=build_risk_items(),
    )

    assert pptx_path.exists(), "工具执行后应生成实际 PPTX 文件。"
    assert "PPT created successfully at:" in result_message, "工具返回消息格式不符合预期。"
    return {"pptx_path": str(pptx_path), "result_message": result_message, "structure_summary": assert_ppt_structure(pptx_path)}


def main() -> int:  # 目的：提供脚本入口；功能：执行映射测试、推导测试和 PPT 冒烟测试并输出 JSON 摘要；实现逻辑：先跑解析级断言，再跑实际导出，最后打印结果和 PASS 标记；可调参数：无；默认参数及原因：返回 0 表示测试通过，方便命令行和 CI 复用。
    assert_alias_mapping()
    assert_derived_metrics()
    assert_financial_values_not_truncated()
    assert_meta_note_from_llm_style_sentence()
    result = run_tool_smoke_test()
    print(json.dumps(result, ensure_ascii=False, indent=2))
    print("PASS")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
